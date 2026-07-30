"""
takeda_finishing.py — Generalized Mode-B FINISHING PASSES (brand-agnostic).

These are the "make it look right" polish passes you run AFTER cloning/recoloring/
retheming a deck, before the QA gate. They encode fixes that recur on every
conversion regardless of source or brand:

  1. align_blade_to_title   — header blade top aligns with the title's first text line
  2. apply_footer_strip_locked / make_strip_bleed
                            — footer as LOCKED layout furniture, behind content,
                              bled edge-to-edge (no gap)
  3. (covered by #2)        — footer behind all slide content => text always in front
  4. fix_text_contrast      — background-aware contrast (white on dark, dark on light),
                              checking slide/layout/master backgrounds AND underlying panels
  5. tidy_title_text        — collapse awkward whitespace, grow the box so a title wraps
                              neatly instead of reproducing a broken source layout

Import alongside takeda:
    import sys; sys.path.insert(0, '/mnt/project')
    import takeda_finishing as fin
    from takeda import _chrome_png            # strip asset source

Then, after your per-slide clone/recolor/retheme and before saving:
    fin.finish_deck(deck.prs, _chrome_png('ENTYVIO_BANNER'),
                    blade_layouts={'Title and Content','Title Only',
                                   'One Third LHS Dark','One Third LHS Light'})

Depends only on python-pptx, lxml, Pillow, numpy (no takeda internals).
"""
import re, io
from pptx.oxml.ns import qn
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
EMU = 914400

# ----------------------------- colour helpers -----------------------------
def _theme_map(prs):
    """Resolve the (possibly rethemed) theme scheme to a {slot: hex} dict, including
    bg1/tx1/bg2/tx2 aliases, so schemeClr references can be resolved to real hexes."""
    m = prs.slide_masters[0]
    for rel in m.part.rels.values():
        if 'theme' in rel.reltype:
            root = etree.fromstring(rel.target_part.blob)
            cs = root.find(f'.//{{{A}}}clrScheme'); d = {}
            for slot in ('dk1','lt1','dk2','lt2','accent1','accent2','accent3',
                         'accent4','accent5','accent6'):
                node = cs.find(f'{{{A}}}{slot}')
                if node is None: continue
                s = node.find(f'{{{A}}}srgbClr'); sy = node.find(f'{{{A}}}sysClr')
                d[slot] = s.get('val') if s is not None else (sy.get('lastClr','000000') if sy is not None else None)
            d['bg1']=d.get('lt1','FFFFFF'); d['tx1']=d.get('dk1','000000')
            d['bg2']=d.get('lt2','FFFFFF'); d['tx2']=d.get('dk2','000000')
            return d
    return {}

def _resolve(el, TH):
    if el is None: return None
    # Accept a bare color container OR a wrapper like <a:rPr> whose color sits
    # inside <a:solidFill>. Searching only direct children of rPr misses every
    # colored run (the 'all accents neutralized to dark' bug).
    for scope in (el, el.find(qn('a:solidFill'))):
        if scope is None: continue
        s = scope.find(qn('a:srgbClr'))
        if s is not None: return s.get('val')
        sc = scope.find(qn('a:schemeClr'))
        if sc is not None: return TH.get(sc.get('val'))
    return None

def _lum(h):
    h = h.lstrip('#'); r,g,b = [int(h[i:i+2],16)/255 for i in (0,2,4)]
    f = lambda c: c/12.92 if c<=0.03928 else ((c+0.055)/1.055)**2.4
    return 0.2126*f(r)+0.7152*f(g)+0.0722*f(b)

def _ratio(a,b):
    la,lb = _lum(a),_lum(b); hi,lo = max(la,lb),min(la,lb)
    return (hi+0.05)/(lo+0.05)

def _shape_fill(sh, TH):
    spPr = sh._element.find(qn('p:spPr'))
    return _resolve(spPr.find(qn('a:solidFill')), TH) if spPr is not None else None

def _base_bg(slide, TH):
    for obj in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        bg = obj.element.find('.//'+qn('p:bg'))
        if bg is not None:
            c = _resolve(bg.find('.//'+qn('a:solidFill')), TH)
            if c: return c
    return 'FFFFFF'

def _box(sh):
    try: return (sh.left/EMU, sh.top/EMU, (sh.left+sh.width)/EMU, (sh.top+sh.height)/EMU)
    except Exception: return None

def _overlap(a,b):
    return not (a[2]<=b[0] or b[2]<=a[0] or a[3]<=b[1] or b[3]<=a[1])

def _set_run_color(rpr, hexv):
    """Set an explicit run/endPara colour, inserting <a:solidFill> in its SCHEMA-CORRECT
    position. In CT_TextCharacterProperties the fill group must come right after <a:ln>
    and BEFORE <a:latin>/<a:ea>/<a:cs>/<a:sym>/<a:hlink*>. Appending it at the end (after
    latin/cs) is tolerated by LibreOffice but IGNORED by PowerPoint — which then falls back
    to the inherited/lstStyle colour (e.g. bg2=purple). This ordering is the fix."""
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        for f in rpr.findall(qn(tag)): rpr.remove(f)
    sf = rpr.makeelement(qn('a:solidFill'), {})
    etree.SubElement(sf, qn('a:srgbClr')).set('val', hexv)
    ln = rpr.find(qn('a:ln'))
    if ln is not None:
        ln.addnext(sf)                 # fill goes immediately after a:ln
    else:
        rpr.insert(0, sf)              # else first child

# ------------------- PASS 4/3: background-aware contrast -------------------
def fix_text_contrast(slide, TH, strip_band=(7.13, 7.5), dark_hex='1A1628', prefer=None):
    """White text on a dark effective background, dark text on a light one.
    Effective bg = the shape's own fill -> the nearest underlying filled shape
    (slide, then layout, then master) -> the footer-strip band -> the slide/layout/
    master background. Colour is written explicitly on runs AND endParaRPr so it
    can't revert to an inherited theme colour. This is the correct fix for
    'headings the same colour as the background' — never blindly force one colour."""
    base = _base_bg(slide, TH)
    shapes = list(slide.shapes)
    for zi, sh in enumerate(shapes):
        if not sh.has_text_frame: continue
        tb = _box(sh)
        eff = _shape_fill(sh, TH)
        if not eff and tb:                                   # underlying slide shapes
            for zj in range(zi-1, -1, -1):
                u = shapes[zj]; ub = _box(u); uf = _shape_fill(u, TH)
                if uf and ub and _overlap(tb, ub): eff = uf; break
        if not eff and tb:                                   # layout, then master shapes
            for container in (slide.slide_layout, slide.slide_layout.slide_master):
                for u in container.shapes:
                    ub = _box(u); uf = _shape_fill(u, TH)
                    if uf and ub and _overlap(tb, ub): eff = uf; break
                if eff: break
        if not eff and tb and strip_band[0] <= tb[1] < strip_band[1]:
            eff = '1A1628'                                   # sits on the dark footer strip
        if not eff: eff = base
        bg_dark = _lum(eff) < 0.5
        target = 'FFFFFF' if bg_dark else dark_hex
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip(): continue
                rpr = run._r.find(qn('a:rPr'))
                if rpr is None:
                    rpr = run._r.makeelement(qn('a:rPr'), {}); run._r.insert(0, rpr)
                cur = _resolve(rpr, TH)
                if bg_dark:
                    if cur != 'FFFFFF': _set_run_color(rpr, 'FFFFFF')
                else:
                    if cur is None:
                        # inherited/uncolored text = NORMAL text (black in the
                        # reference) — stamp neutral dark, NEVER a brand accent.
                        _set_run_color(rpr, target)
                    elif _ratio(cur, eff) >= 3.0:
                        pass                     # readable accent — PRESERVE it
                    else:
                        # the text WAS deliberately colored but now fails contrast:
                        # correct WITHIN the brand family (preference order) first
                        repl = target
                        if prefer:
                            for cand in prefer:
                                if cand and _ratio(cand, eff) >= 3.0:
                                    repl = cand; break
                        _set_run_color(rpr, repl)
            epr = para._p.find(qn('a:endParaRPr'))
            if epr is not None: _set_run_color(epr, target)

# --------------------- PASS 5: tidy title whitespace/fit ---------------------
def tidy_title_text(slide, max_h_in=2.2):
    """Collapse runs of spaces, strip spaces around soft line-breaks / line ends, and
    grow the title box (height only) so it never clips — so an awkward source title
    wraps neatly instead of being reproduced verbatim. NEVER writes a partial xfrm
    (a bare <a:off> without y collapses the box to zero width in PowerPoint), so it
    only resizes when a complete off+ext already exists."""
    for sh in slide.shapes:
        if not (sh.is_placeholder and sh.has_text_frame): continue
        try:
            if sh.placeholder_format.type is not None and sh.placeholder_format.idx != 0:
                continue
        except Exception:
            continue
        tb = sh.text_frame
        for para in tb.paragraphs:
            runs = para.runs
            for k, run in enumerate(runs):
                t = re.sub(r' {2,}', ' ', run.text)
                nxt = run._r.getnext()
                if k == len(runs)-1 or (nxt is not None and nxt.tag == qn('a:br')):
                    t = t.rstrip()
                prv = run._r.getprevious()
                if prv is not None and prv.tag == qn('a:br'):
                    t = t.lstrip()
                if t != run.text: run.text = t
        n_lines = 0; max_sz = 24
        for para in tb.paragraphs:
            n_lines += 1 + len(para._p.findall(qn('a:br')))
            for run in para.runs:
                if run.font.size: max_sz = max(max_sz, run.font.size.pt)
        need = n_lines * (max_sz * 1.2 / 72.0) + 0.15
        spPr = sh._element.find(qn('p:spPr'))
        if spPr is None: continue
        xfrm = spPr.find(qn('a:xfrm'))
        if xfrm is None: continue
        off = xfrm.find(qn('a:off')); ext = xfrm.find(qn('a:ext'))
        if off is None or ext is None: continue
        if None in (off.get('x'), off.get('y'), ext.get('cx'), ext.get('cy')): continue
        if need > int(ext.get('cy'))/EMU and need <= max_h_in:
            ext.set('cy', str(int(need*EMU)))                # width & position preserved

# ---------------------- PASS 1: align blade to title ----------------------
def align_blade_to_title(slide, blade_name='TakedaBlade'):
    """Move the header blade so its TOP edge meets the title's first text-line top
    (title box top + top inset), instead of floating at a fixed Y."""
    blade = title = None
    for sh in slide.shapes:
        if getattr(sh, 'name', '') == blade_name: blade = sh
        if sh.is_placeholder and sh.has_text_frame:
            try:
                if sh.placeholder_format.idx == 0: title = sh
            except Exception: pass
    if blade is None or title is None or title.top is None: return
    bodyPr = title._element.find('.//'+qn('a:bodyPr'))
    t_ins = int(bodyPr.get('tIns')) if (bodyPr is not None and bodyPr.get('tIns')) else 45720
    off = blade._element.find('.//'+qn('a:xfrm')).find(qn('a:off'))
    off.set('y', str(int(title.top) + t_ins))

# ---------------- PASS 2/3: locked, edge-to-edge footer strip ----------------
def make_strip_bleed(src_png, out_png='/tmp/strip_bleed.png', sat_thresh=0.30):
    """Crop the near-white margins baked into the footer asset's edges so the
    saturated gradient bleeds edge-to-edge (fixes the perimeter 'gap')."""
    from PIL import Image
    import numpy as np
    im = Image.open(src_png).convert('RGB')
    a = np.asarray(im).astype(float)
    mx = a.max(axis=2); mn = a.min(axis=2)
    sat = np.where(mx > 0, (mx - mn) / np.clip(mx, 1, None), 0)
    content = sat > sat_thresh
    cols = np.where(content.mean(axis=0) > 0.4)[0]
    rows = np.where(content.mean(axis=1) > 0.4)[0]
    if len(cols) == 0 or len(rows) == 0:
        im.save(out_png); return out_png
    im.crop((int(cols.min()), int(rows.min()),
             int(cols.max())+1, int(rows.max())+1)).save(out_png)
    return out_png

def apply_footer_strip_locked(prs, strip_png, band=(0.0, 7.15, 13.333, 0.35),
                              bleed=True, name='entyvio_footer_strip'):
    """Place the footer strip once on every USED layout as LOCKED, non-selectable
    furniture: it renders behind all slide content (so text is always in front) and
    above each layout's own background (so it shows even on dark dividers), and it
    cannot be clicked/moved/resized on a slide. Removes any per-slide strip pics first.
    bleed=True crops the asset margins so the gradient reaches both edges."""
    if bleed: strip_png = make_strip_bleed(strip_png)
    L, T, W, H = band
    for s in prs.slides:
        for sh in list(s.shapes):
            if sh.shape_type == 13 and sh.top is not None and abs(sh.top/EMU - T) < 0.06:
                sh._element.getparent().remove(sh._element)
    seen, used = set(), []
    for s in prs.slides:
        lay = s.slide_layout
        if id(lay) not in seen: seen.add(id(lay)); used.append(lay)
    with open(strip_png, 'rb') as f: blob = f.read()
    for lay in used:
        if any(getattr(sh,'name','')==name for sh in lay.shapes): continue
        _img, rId = lay.part.get_or_add_image_part(io.BytesIO(blob))
        pic = f'''<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">
          <p:nvPicPr>
            <p:cNvPr id="990" name="{name}"/>
            <p:cNvPicPr><a:picLocks noSelect="1" noMove="1" noResize="1" noChangeAspect="1"/></p:cNvPicPr>
            <p:nvPr userDrawn="1"/>
          </p:nvPicPr>
          <p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
          <p:spPr><a:xfrm><a:off x="{int(L*EMU)}" y="{int(T*EMU)}"/>
            <a:ext cx="{int(W*EMU)}" cy="{int(H*EMU)}"/></a:xfrm>
            <a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
        </p:pic>'''
        lay.shapes._spTree.insert(2, etree.fromstring(pic))   # behind content, above layout bg

# ------------------------------ orchestrator ------------------------------
def finish_deck(prs, strip_png, blade_layouts=None, add_blade=None):
    """Run every finishing pass in the correct order. `add_blade` is an optional
    callable add_blade(slide) (e.g. takeda.add_blade_accent) used on slides whose
    layout name is in `blade_layouts`. Contrast runs LAST so titles/headings are
    coloured against their real, final backgrounds."""
    TH = _theme_map(prs)
    for slide in prs.slides:
        tidy_title_text(slide)
    if add_blade and blade_layouts:
        for slide in prs.slides:
            if slide.slide_layout.name in blade_layouts:
                add_blade(slide)
                align_blade_to_title(slide)
    apply_footer_strip_locked(prs, strip_png)
    for slide in prs.slides:
        fix_text_contrast(slide, TH)


# ═══════════════════════════════════════════════════════════════════════════════
# PATH A — IN-PLACE RESKIN (added 2026-07-29)
# Mode B DEFAULT. Restyles the SOURCE package in place: retheme → recolor →
# refont → chrome → title style → contrast. No new Deck, no cross-package shape
# cloning, no part grafting. Rationale: template migration (fresh Takeda Deck +
# clone_shapes_with_rels) re-serializes the 5-master corporate template and
# produces packages PowerPoint's strict reader "repairs" (strips content on
# open) even though python-pptx / LibreOffice / XSD all accept them. The
# in-place reskin keeps the container PowerPoint already trusts, so repair is
# impossible by construction. Verified: ZS COGNITIVE 36-slide → ENTYVIO opens
# clean in PowerPoint.
# ═══════════════════════════════════════════════════════════════════════════════

TITLE_PT = 25            # org rule (2026-07-29): headers are ALWAYS 25pt, all formats
TITLE_FONT = 'Aptos'     # org rule: font is always Aptos

_NON_CONTENT = ('Divider', 'Title Page')      # substring matches (divider/cover families)
_NON_CONTENT_EXACT = ('Title', 'End', 'End Slide', 'Cover')   # exact layout names


def _is_display_layout(name):
    """Dividers, covers and end slides carry display titles, not headers."""
    return name in _NON_CONTENT_EXACT or any(k in name for k in _NON_CONTENT)


def enforce_title_style(prs, size_pt=TITLE_PT, font=TITLE_FONT, safe=None):
    """Force every content-slide title to the org standard (25pt Aptos, all brands).

    Source decks usually leave titles inheriting the master titleStyle (e.g. ZS
    master = 32pt) with PowerPoint autofit silently shrinking them, so 'preserve
    source sizes' produces surprising header sizes. This pass: strips stale
    normAutofit scaling, writes explicit sz/typeface on every run (adding rPr
    where missing), top-anchors the box, grows its height for the new wrap count
    (complete xfrm), and re-aligns the blade. Skips dividers/cover/end slides —
    those are display titles, not headers."""
    import math
    from pptx.util import Inches, Emu
    from lxml import etree
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    n_done = 0
    for slide in prs.slides:
        lname = slide.slide_layout.name
        if _is_display_layout(lname):
            continue
        for sh in slide.shapes:
            if not (sh.is_placeholder and sh.has_text_frame):
                continue
            try:
                if sh.placeholder_format.idx != 0:
                    continue
            except Exception:
                continue
            txBody = sh.text_frame._txBody
            bodyPr = txBody.find(f'{{{A}}}bodyPr')
            if bodyPr is None:
                continue
            for auto in bodyPr.findall(f'{{{A}}}normAutofit'):
                bodyPr.remove(auto)              # stale fontScale would defeat sz
            panel_layout = 'One Third' in lname  # panel titles are DESIGNED centered —
            if not panel_layout:                 # keep their anchor/box; size+font only
                bodyPr.set('anchor', 't')
            for tag in ('rPr', 'endParaRPr'):
                for r in txBody.iter(f'{{{A}}}{tag}'):
                    r.set('sz', str(size_pt * 100))
                    latin = r.find(f'{{{A}}}latin')
                    if latin is None:
                        latin = etree.SubElement(r, f'{{{A}}}latin')
                    latin.set('typeface', font)
            for para in txBody.iter(f'{{{A}}}p'):
                for run in para.findall(f'{{{A}}}r'):
                    if run.find(f'{{{A}}}rPr') is None:
                        rPr = etree.SubElement(run, f'{{{A}}}rPr')
                        rPr.set('sz', str(size_pt * 100))
                        latin = etree.SubElement(rPr, f'{{{A}}}latin')
                        latin.set('typeface', font)
                        run.insert(0, rPr)
            # grow the box for the new wrap count — write a COMPLETE xfrm.
            # ⚠️ Placeholders often INHERIT geometry from the layout (no xfrm on
            # the slide element). python-pptx setters then create a fresh xfrm
            # containing only what you set — writing height alone yields
            # <a:ext cx="0" .../> and the title collapses to a zero-width sliver.
            # Always read the (possibly inherited) L/T/W/H and write ALL FOUR.
            L, T, W, H = sh.left, sh.top, sh.width, sh.height
            if not panel_layout and None not in (L, T, W, H) and W > 0:
                # brand title SAFE ZONE: cap the right edge so the title never
                # collides with top-banner chrome (e.g. FRUQ's corner wave)
                if safe and safe.get('max_right_in'):
                    max_r = Inches(safe['max_right_in'])
                    if L + W > max_r and max_r > L:
                        W = max_r - L
                w = Emu(W).inches
                cpl = max(20, int(w * 152 / size_pt))
                lines = max(1, math.ceil(len(sh.text_frame.text) / cpl))
                need = min(0.017 * size_pt * lines + 0.14, 1.8)
                sh.left, sh.top, sh.width = L, T, W
                sh.height = H if Emu(H).inches >= need else Inches(need)
            align_blade_to_title(slide)
            n_done += 1
            break
    return n_done


def remove_offcanvas_shapes(slide):
    """Delete shapes lying entirely left of the canvas (parked source notes).
    They are invisible junk that trips the QA left-zone check."""
    from pptx.util import Emu
    removed = 0
    for sh in list(slide.shapes):
        if sh.left is not None and sh.width is not None \
           and Emu(sh.left).inches + Emu(sh.width).inches <= 0:
            sh._element.getparent().remove(sh._element)
            removed += 1
    return removed


def grow_overflow_bars(slide):
    """Grow (height only) dark filled text bars whose content no longer fits —
    Aptos runs wider than Arial at equal size, so refonted description bars can
    crop their last line. Conservative estimate; complete xfrm preserved."""
    import math
    from pptx.util import Inches, Emu
    grown = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            if sh.fill.type != 1:
                continue
            r, g, b = sh.fill.fore_color.rgb
            if (r + g + b) >= 300:
                continue
        except Exception:
            continue
        tf = sh.text_frame
        if not tf.text.strip():
            continue
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        if w < 3 or h > 2.5:
            continue
        pts = [r_.font.size.pt for p in tf.paragraphs for r_ in p.runs if r_.font.size]
        pt = max(pts) if pts else 11
        cpl = max(10, int(w * 150 / pt))
        lines = sum(max(1, math.ceil(len(p.text) / cpl)) for p in tf.paragraphs)
        need = lines * pt * 1.28 / 72 + 0.12
        if need > h + 0.02:
            sh.height = Inches(min(need, 2.5))
            grown += 1
    return grown


def reskin_deck(src_path, out_path, brand, extra_remap=None,
                title_pt=TITLE_PT, title_font=TITLE_FONT):
    """MODE B, PATH A — one-call in-place reskin of a source deck to a Takeda brand.

        import takeda_finishing as fin
        fin.reskin_deck('src.pptx', 'out.pptx', 'ENTYVIO')

    Sequence (order matters): retheme → recolor+refont masters/layouts →
    per-slide chrome strip-out, junk removal, recolor, refont → overflow bars →
    locked footer strip + blades → title style (25pt Aptos) → contrast LAST →
    save → clean_pptx_zip. Then run the QA gate with --reformat."""
    import takeda
    from pptx import Presentation

    b = takeda.get_brand(brand)
    remap = dict(b['BRAND_REMAP'])
    if extra_remap:
        remap.update(extra_remap)
    theme_key = next((k for k in b if k.endswith('_THEME_REMAP')), None)

    prs = Presentation(src_path)
    if theme_key:
        takeda.retheme_brand(prs, b[theme_key])

    seen = set()
    for slide in prs.slides:
        lay = slide.slide_layout
        for part in (lay, lay.slide_master):
            key = str(part.part.partname)
            if key in seen:
                continue
            seen.add(key)
            takeda.recolor_element(part.shapes._spTree, remap)
            takeda.refont_element(part.shapes._spTree, title_font)

    for slide in prs.slides:
        takeda.strip_source_chrome(slide)
        remove_offcanvas_shapes(slide)
        takeda.recolor_element(slide.shapes._spTree, remap)
        takeda.refont_element(slide.shapes._spTree, title_font)
        grow_overflow_bars(slide)

    # ── brand chrome (heterogeneous per brand — see each brand-*.md) ──
    if brand == 'ICLUSIG' and hasattr(takeda, 'stamp_iclusig_chrome'):
        takeda.stamp_iclusig_chrome(prs)          # deck-level: layouts, dividers, logo
    elif brand == 'CORPORATE':
        for slide in prs.slides:
            lname = slide.slide_layout.name
            if _is_display_layout(lname) or 'One Third' in lname or 'LHS' in lname:
                continue          # panel layouts carry their own composition — no blade
            takeda.add_blade_accent(slide)
            align_blade_to_title(slide)
    elif brand == 'ENTYVIO':
        try:
            strip_png = takeda._chrome_png('ENTYVIO_BANNER')
        except Exception:
            strip_png = None
        if strip_png:
            apply_footer_strip_locked(prs, strip_png)
        for slide in prs.slides:
            lname = slide.slide_layout.name
            if _is_display_layout(lname) or 'One Third' in lname or 'LHS' in lname:
                continue          # panel layouts carry their own composition — no blade
            takeda.add_blade_accent(slide)
            align_blade_to_title(slide)
    else:
        # FRUZAQLA / ALUNBRIG-style per-slide chrome via the dispatcher
        # (FRUQ: top gradient band behind content + whitespace-contingent logo).
        for slide in prs.slides:
            lname = slide.slide_layout.name
            # 'One Third' panel layouts carry their own full-height side-panel
            # composition — a top band there covers the panel head / titles.
            if _is_display_layout(lname) or 'One Third' in lname:
                continue
            try:
                takeda.apply_brand_chrome(slide, brand)
            except Exception:
                pass

    # ── divider backgrounds: brand identity color (before contrast, which whitens titles) ──
    div_bg = b.get('BRAND_DIVIDER_BG')
    if div_bg:
        from pptx.dml.color import RGBColor
        rgb = RGBColor(int(div_bg[0:2], 16), int(div_bg[2:4], 16), int(div_bg[4:6], 16))
        for slide in prs.slides:
            if 'Divider' in slide.slide_layout.name:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = rgb
        recolor_side_panels(prs, div_bg)   # dark side panels, SHAPE-level (never via
        #                                    the accent1 theme slot — that slot is TEXT)

    enforce_title_style(prs, title_pt, title_font, safe=b.get('BRAND_TITLE_SAFE'))

    TH = _theme_map(prs)
    prefer = b.get('BRAND_TEXT_ACCENTS')
    dark_hex = b.get('BRAND_DARK_TEXT', '1A1628')
    for slide in prs.slides:
        fix_text_contrast(slide, TH, dark_hex=dark_hex, prefer=prefer)

    icon_map = b.get('BRAND_ICON_RECOLOR')
    if icon_map:
        recolor_svg_icons(prs, icon_map)

    apply_brand_cover(prs, b)   # after contrast: cover text whites stay explicit

    prs.save(out_path)
    takeda.clean_pptx_zip(out_path)
    return out_path


def recolor_svg_icons(prs, hex_map):
    """Recolor SVG icon PARTS in place (source line icons are SVG pictures —
    shape-level remaps cannot reach them). Rewrites color literals inside each
    .svg part's XML; content-only edit, so no package-structure/repair risk.
    Note: the PNG raster fallbacks keep their source color (only pre-2019
    PowerPoint uses them)."""
    changed = 0
    for part in prs.part.package.iter_parts():
        if str(part.partname).lower().endswith('.svg'):
            try:
                body = part.blob.decode('utf-8')
            except Exception:
                continue
            new = body
            for old_hex, new_hex in hex_map.items():
                for variant in (old_hex.upper(), old_hex.lower()):
                    new = new.replace('#' + variant, '#' + new_hex.upper())
            if new != body:
                part._blob = new.encode('utf-8')
                changed += 1
    return changed


def recolor_side_panels(prs, hex_bg):
    """Recolor tall dark side panels (e.g. ZS 'One Third LHS Dark') to the brand
    divider color at SHAPE level. ZS fills these with schemeClr accent1 — its
    near-black TEXT slot — so the slot itself must never be remapped to a
    chromatic (it recolors all body text). Fix the panel shapes directly."""
    from pptx.util import Emu
    from lxml import etree
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    seen, n = set(), 0
    def fix(shapes):
        nonlocal n
        for sh in shapes:
            if sh.width is None or sh.height is None: continue
            if Emu(sh.width).inches < 2 or Emu(sh.height).inches < 4.5: continue
            # NB: shape spPr lives in the PRESENTATIONML namespace (p:spPr) — an
            # a:-namespace search finds nothing and fails silently. Use oxml's accessor.
            spPr = getattr(sh._element, 'spPr', None)
            if spPr is None: continue
            sf = spPr.find(f'{{{A}}}solidFill')
            if sf is None: continue
            sc = sf.find(f'{{{A}}}schemeClr'); srgb = sf.find(f'{{{A}}}srgbClr')
            is_dark_panel = (sc is not None and sc.get('val') in ('accent1', 'tx2', 'dk2')) or \
                            (srgb is not None and sum(int(srgb.get('val')[i:i+2], 16) for i in (0,2,4)) < 220)
            if is_dark_panel:
                for child in list(sf): sf.remove(child)
                etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', hex_bg)
                n += 1
    for slide in prs.slides:
        lname = slide.slide_layout.name
        if 'One Third' in lname or 'LHS' in lname:
            fix(slide.shapes)
            key = str(slide.slide_layout.part.partname)
            if key not in seen:
                seen.add(key)
                fix(slide.slide_layout.shapes)
    return n


def apply_brand_cover(prs, brand_dict):
    """Stamp the brand's full-bleed cover art on the COVER slide (layout 'Title')
    and whiten its text. In-place & PowerPoint-safe: one image part registered via
    get_or_add_image_part + slide-level shape edits — no cross-package grafting.
    The art file resolves local-first, then from the GitHub repo (takeda_remote),
    so dropping real brand artwork into the repo overrides the procedural default."""
    asset = brand_dict.get('BRAND_COVER_ASSET')
    logos = brand_dict.get('BRAND_COVER_LOGOS')
    if not asset and not logos:
        return False
    import io, copy
    import takeda_remote
    from pptx.util import Emu
    from lxml import etree
    art_path = None
    if asset:
        try:
            art_path = takeda_remote.get(asset)
        except Exception:
            art_path = None
    if not art_path and not logos:
        return False
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    done = False
    for slide in prs.slides:
        if slide.slide_layout.name not in ('Title', 'Title Page 1', 'Title Slide'):
            continue
        spTree = slide.shapes._spTree
        if art_path or logos:
            # drop the source's decorative pictures/logo art (they fight the new
            # art AND the stamped logos — e.g. the ZS mark under the Takeda logo)
            for sh in list(slide.shapes):
                if sh.shape_type == 13:
                    sh._element.getparent().remove(sh._element)
            # source cover logos often live on the LAYOUT (render beneath the
            # slide and show through when there is no full-bleed art) — clear
            # layout pictures that intersect any stamped logo zone
            if logos:
                from pptx.util import Emu as _Emu
                zones = [(x, y, x + w, y + w) for _, x, y, w in logos]
                for lsh in list(slide.slide_layout.shapes):
                    if lsh.shape_type != 13 or lsh.left is None or lsh.width is None:
                        continue
                    L0, T0 = _Emu(lsh.left).inches, _Emu(lsh.top).inches
                    R0, B0 = L0 + _Emu(lsh.width).inches, T0 + _Emu(lsh.height).inches
                    if any(L0 < zx2 + 0.3 and R0 > zx1 - 0.3 and T0 < zy2 + 0.3 and B0 > zy1 - 0.3
                           for zx1, zy1, zx2, zy2 in zones):
                        lsh._element.getparent().remove(lsh._element)
        if art_path:
            with open(art_path, 'rb') as f:
                image_part, rId = slide.part.get_or_add_image_part(io.BytesIO(f.read()))
        cx, cy = prs.slide_width, prs.slide_height
        if art_path:
            pic_xml = (
                f'<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">'
                f'<p:nvPicPr><p:cNvPr id="991" name="brand_cover_art"/>'
                f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
                f'<p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
                f'<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
            from pptx.oxml import parse_xml         # oxml classes, not raw lxml —
            pic = parse_xml(pic_xml)                # raw elements break shape iteration
            spTree.insert(2, pic)                   # behind all content
        # brand + Takeda logos in the artwork's reserved strip
        for asset, x_in, y_in, w_in in brand_dict.get('BRAND_COVER_LOGOS', []):
            try:
                import takeda
                path = None
                try:
                    path = takeda._chrome_png(asset)      # registry asset (e.g. ICLUSIG_LOGO)
                except Exception:
                    path = takeda_remote.get(asset)       # repo/local file (e.g. takeda_logo.png)
                if not path: continue
                with open(path, 'rb') as lf:
                    lp, lrId = slide.part.get_or_add_image_part(io.BytesIO(lf.read()))
                from PIL import Image as _Im
                iw, ih = _Im.open(path).size
                w_emu = int(w_in * 914400); h_emu = int(w_emu * ih / iw)
                lx = f'<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">' \
                     f'<p:nvPicPr><p:cNvPr id="99{hash(asset) % 90 + 2}" name="cover_logo_{asset[:12]}"/>' \
                     f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>' \
                     f'<p:blipFill><a:blip r:embed="{lrId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>' \
                     f'<p:spPr><a:xfrm><a:off x="{int(x_in*914400)}" y="{int(y_in*914400)}"/>' \
                     f'<a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>' \
                     f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
                from pptx.oxml import parse_xml
                spTree.append(parse_xml(lx))              # on top of the art
            except Exception:
                pass
        if not art_path:
            done = True
            continue                               # logos only — keep source text colors
        ct = brand_dict.get('BRAND_COVER_TEXT', {})
        title_hex = ct.get('title', 'FFFFFF')
        body_hex = ct.get('body', 'FFFFFF')      # default: dark art -> all white
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                is_title = False
                try:
                    is_title = sh.is_placeholder and sh.placeholder_format.idx == 0
                except Exception:
                    pass
                if not is_title:
                    is_title = 'title' in (sh.name or '').lower() or 'pres_nam' in (sh.name or '').lower()
                hexval = title_hex if is_title else body_hex
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        _set_run_color(run._r.get_or_add_rPr(), hexval)
        done = True
    return done


def apply_logo_lockup(slide, logos, left=0.55, top=0.35, height=0.72,
                      gap=0.2, divider_hex='A1A4AC'):
    """Arrange multiple logos side by side as a lockup — 'LogoA | LogoB' — with
    thin divider bars between, all vertically centered on one axis.

    THE LOGO CLASH RULE: when two logos compete for the same zone (a stamped
    brand logo over a source/partner logo, or any co-branding), never overlap
    them and never delete one arbitrarily — build a lockup. WHICH logos to keep
    is a judgment/user call; this helper only does the arranging.

        fin.apply_logo_lockup(slide, ['takeda_logo.png', 'zs_logo.png'])
        fin.apply_logo_lockup(slide, ['ICLUSIG_LOGO', '/path/to/partner.png'])

    `logos`: image paths, repo filenames (takeda_remote), or chrome asset names.
    Existing pictures inside the lockup zone are cleared first. Returns the
    number of logos placed."""
    import io
    import takeda
    import takeda_remote
    from PIL import Image as _Im
    from pptx.util import Emu
    from pptx.oxml import parse_xml
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    resolved = []
    for item in logos:
        path = None
        try:
            path = takeda._chrome_png(item)
        except Exception:
            try:
                path = takeda_remote.get(item)
            except Exception:
                path = item if isinstance(item, str) else None
        if path:
            try:
                iw, ih = _Im.open(path).size
                resolved.append((path, iw, ih))
            except Exception:
                pass
    if len(resolved) < 1:
        return 0

    widths = [height * iw / ih for _, iw, ih in resolved]
    total = sum(widths) + (len(resolved) - 1) * (gap * 2 + 0.012)
    # clear existing pictures in the lockup zone
    zx1, zy1, zx2, zy2 = left - 0.1, top - 0.1, left + total + 0.1, top + height + 0.1
    for coll in (slide.shapes, slide.slide_layout.shapes):
        for sh in list(coll):
            if sh.shape_type != 13 or sh.left is None or sh.width is None:
                continue
            L0, T0 = Emu(sh.left).inches, Emu(sh.top).inches
            R0, B0 = L0 + Emu(sh.width).inches, T0 + Emu(sh.height).inches
            if L0 < zx2 and R0 > zx1 and T0 < zy2 and B0 > zy1:
                sh._element.getparent().remove(sh._element)

    x = left
    spTree = slide.shapes._spTree
    for i, ((path, iw, ih), w) in enumerate(zip(resolved, widths)):
        with open(path, 'rb') as f:
            _, rId = slide.part.get_or_add_image_part(io.BytesIO(f.read()))
        pic = parse_xml(
            f'<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">'
            f'<p:nvPicPr><p:cNvPr id="98{i+1}" name="lockup_logo_{i+1}"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{int(x*914400)}" y="{int(top*914400)}"/>'
            f'<a:ext cx="{int(w*914400)}" cy="{int(height*914400)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
        spTree.append(pic)
        x += w
        if i < len(resolved) - 1:
            bar = parse_xml(
                f'<p:sp xmlns:p="{P}" xmlns:a="{A}">'
                f'<p:nvSpPr><p:cNvPr id="97{i+1}" name="lockup_divider_{i+1}"/>'
                f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="{int((x+gap)*914400)}" y="{int((top+0.05)*914400)}"/>'
                f'<a:ext cx="{int(0.012*914400)}" cy="{int((height-0.1)*914400)}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'<a:solidFill><a:srgbClr val="{divider_hex}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></p:spPr>'
                f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')
            spTree.append(bar)
            x += gap * 2 + 0.012
    return len(resolved)
