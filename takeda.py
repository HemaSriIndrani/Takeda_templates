"""
takeda.py — Takeda deck generation toolkit (single-file edition).

Contains:
  • Deck class + layout registry + palette/contrast helpers  (Mode A)
  • Mode B helpers: clone, recolor, chrome, retheme          (takeda_modeb)
  • Per-brand color registry: get_brand("ICLUSIG")           (takeda_brands)
  • Executive-summary slide builders                          (takeda_execsum)
  • QA gate: run as  python takeda.py output.pptx            (validate)

Usage — import:
    from takeda import Deck, get_brand, clone_shapes_with_rels, clean_pptx_zip

Usage — validate:
    python takeda.py output.pptx
    # exit 0 = pass (safe to share), exit 1 = fail (fix and re-run)

Usage — regenerate layout map after template change:
    python takeda.py --regen-layouts Takeda_Slide_Template_EN.potx
"""

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — CORE DECK API (Deck class, layout registry, palette, helpers)
# ═══════════════════════════════════════════════════════════════════════════════
"""
takeda_deck.py — Takeda 2026 deck-generation library.
Auto-extracted from TAKEDA-TEMPLATE-API.md (do not hand-edit; regenerate from source).
Provides: Deck class, layout registry, Mode B helpers, palette audit, contrast,
blade accent, chart recoloring, relationship-aware shape cloning, zip cleanup.
"""

# ===== Block 1: Deck API, constants, layout registry, module helpers =====
"""
Takeda Deck API — wraps Takeda_Slide_Template_EN.potx in a clean Python API.
Handles .potx → .pptx content-type conversion automatically.
"""
import io, os, zipfile, copy, math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
_MST2 = _MST  # alias used by icon-centering helper
_T = _MST     # alias used by audit_mode_b_fidelity
import lxml.etree as etree

# ── BRAND CONSTANTS ──────────────────────────────────────────
# Updated for the 2026 Takeda template refresh (Aptos font, new neutral accents).
# accent1/2/3 (red / dark-red / charcoal) are UNCHANGED from the prior template.
TAKEDA_RED    = RGBColor(0xE1, 0x24, 0x2A)   # accent1 — unchanged
DARK_RED      = RGBColor(0x89, 0x15, 0x15)   # accent2 — unchanged
CHARCOAL      = RGBColor(0x34, 0x37, 0x3F)   # accent3 / dk1 — unchanged
MID_GRAY      = RGBColor(0xA1, 0xB1, 0xC3)   # accent4 — CHANGED (was A1A4AC; now blue-gray)
LIGHT_GRAY    = RGBColor(0xD1, 0xD8, 0xE0)   # accent5 — CHANGED (was EDF2F4; now cooler)
PALE_GRAY     = RGBColor(0xED, 0xF2, 0xF3)   # accent6 — CHANGED (was FFFFFF; now pale blue-gray)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)
BRAND_FONT    = "Aptos"          # body font — CHANGED from Calibri (2026 refresh)
BRAND_FONT_DISPLAY = "Aptos Display"  # title/display font

# Palette audit set — includes BOTH the new neutrals AND the legacy ones, because
# Mode B source decks and older Takeda decks still carry the legacy grays. Treat
# legacy neutrals as acceptable (they're near-identical perceptually).
TAKEDA_PALETTE = {
    'E1242A','891515','34373F',                 # core accents (unchanged)
    'A1B1C3','D1D8E0','EDF2F3',                  # NEW neutrals (2026)
    'A1A4AC','EDF2F4',                           # legacy neutrals (still tolerated)
    'FFFFFF','000000',
}

# ── LAYOUT CATALOG ───────────────────────────────────────────
# Updated for the 2026 template (Master 0 = 24 layouts). Indices below are
# master-0 positions as python-pptx sees them (prs.slide_layouts). Names are the
# source of truth — get_layout() resolves by NAME, so a future re-order won't break.
# Four extra "Pattern" masters (1-4) carry full-bleed branded backgrounds for
# title/divider slides; see PATTERN_LAYOUTS below.
LAYOUTS = {
    # Title / opener variants (all now ship with Pattern backgrounds baked into the layout)
    "TITLE_PAGE_PHOTO":          (0,  "Title Page 1"),       # hero photo + pattern
    "TITLE_PAGE_4":              (1,  "4_Title Page 1"),
    "TITLE_PAGE_3":              (2,  "3_Title Page 1"),
    "TITLE_PAGE_NO_PHOTO":       (3,  "1_Title Page 1"),     # **Default opener** (no photo)
    # Section dividers — NEW: three pattern-background divider variants (replaces old dark panel)
    "DIVIDER_2":                 (4,  "2_Divider 1"),
    "DIVIDER_1":                 (5,  "1_Divider 1"),         # **Default divider**
    "DIVIDER_3":                 (6,  "3_Divider 1"),
    # Content
    "ONE_COLUMN":                (7,  "Standard 1-Column Text"),    # **Default content**
    "ONE_COLUMN_WAVE":           (8,  "1_Standard 1-Column Text"),
    "END_SLIDE":                 (9,  "End Slide"),                 # **Always end with this**
    "BLANK":                     (10, "Blank slide"),
    "TWO_COLUMN_ADVANCED":       (11, "Advanced 2-Column Text"),
    "THREE_COLUMN":              (12, "Advanced 3-Column Text"),
    "FOUR_COLUMN":               (13, "Advanced 4-Column Text"),
    "IMAGE_TWO_THIRDS":          (14, "Advanced 2/3 Image"),
    "IMAGE_HALF":                (15, "Advanced _1/2 image"),
    "IMAGE_HALF_FULLBLEED":      (16, "Advanced 1/2 Image (Full Bleed)"),
    "IMAGE_ONE_THIRD":           (17, "Advanced 1/3 Image"),
    "IMAGE_ONE_THIRD_FULLBLEED": (18, "Advanced 1/3 Image (Full Bleed)"),
    "IMAGE_FOUR_COLUMN":         (19, "Advanced Image 4-Column"),
    "CHART_FULL":                (20, "Advanced Chart Full Width"),
    "CHART_TWO_THIRDS":          (21, "Advanced Chart 2/3"),
    "CHART_TWO_COLUMN":          (22, "Advanced Chart 2 Column"),
    "BIG_PICTURE":               (23, "Big Picture"),               # NEW full-bleed image layout
}

# NEW in 2026: four "Pattern" masters, each with 4 background treatments. These are
# decorative full-bleed branded backgrounds (red flowing-line + scallop motifs).
# POLICY: use Pattern backgrounds for TITLE and DIVIDER slides ONLY. Content slides
# stay clean (white background) for readability. Resolve via get_pattern_layout().
PATTERN_LAYOUTS = {
    # (master_index, layout_index_within_master): description
    "PATTERN_1_GRADIENT":  (1, 0), "PATTERN_1_INVERTED": (1, 1),
    "PATTERN_1_SUBTLE_1":  (1, 2), "PATTERN_1_SUBTLE_2": (1, 3),
    "PATTERN_2_GRADIENT":  (2, 0), "PATTERN_2_INVERTED": (2, 1),
    "PATTERN_2_SUBTLE_1":  (2, 2), "PATTERN_2_SUBTLE_2": (2, 3),
    "PATTERN_3_GRADIENT":  (3, 0), "PATTERN_3_INVERTED": (3, 1),
    "PATTERN_3_SUBTLE_1":  (3, 2), "PATTERN_3_SUBTLE_2": (3, 3),
    "PATTERN_4_GRADIENT":  (4, 0), "PATTERN_4_INVERTED": (4, 1),
    "PATTERN_4_SUBTLE_1":  (4, 2), "PATTERN_4_SUBTLE_2": (4, 3),
}

# Removed in 2026 (present in the old template, gone now). If a caller references
# one of these, fall back to the mapped replacement.
RETIRED_LAYOUTS = {
    "TITLE_PAGE_RIGHT_PHOTO": "TITLE_PAGE_PHOTO",   # "Title Page 2" removed
    "TWO_COLUMN":             "TWO_COLUMN_ADVANCED", # "Standard 2-Column Text" removed
    "SECTION_TITLE":          "DIVIDER_1",           # "Section Title Page" → new dividers
    "CHART_THREE_COLUMN":     "CHART_TWO_COLUMN",    # "Advanced Chart 3 Column" removed
    "CONTENT_SLIDE":          "ONE_COLUMN",          # "Slide of content" removed
    "TITLE_AND_CONTENT":      "ONE_COLUMN",          # "Title and Content" removed
    "TITLE_ONLY":             "ONE_COLUMN",          # "Title Only" removed
    "TITLE_ONLY_JP":          "ONE_COLUMN",          # locale dup removed
    "TITLE_PAGE_PHOTO_ALT":   "TITLE_PAGE_4",        # "1_Title Page 1" repurposed
}

# ── .potx → .pptx CONVERSION ────────────────────────────────
_TEMPLATE_CT = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_PRESENTATION_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

def _potx_to_pptx_bytes(potx_path):
    src_bytes = Path(potx_path).read_bytes()
    buf_in, buf_out = io.BytesIO(src_bytes), io.BytesIO()
    with zipfile.ZipFile(buf_in, "r") as zin, zipfile.ZipFile(buf_out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(_TEMPLATE_CT.encode(), _PRESENTATION_CT.encode())
            zout.writestr(item, data)
    return buf_out.getvalue()

def _find_template(brand="CORPORATE"):
    """Locate the template file for a brand, on disk or from takeda_templates.

    `brand` ∈ {CORPORATE, FRUZAQLA, ICLUSIG, ENTYVIO}. CORPORATE is the default
    Takeda 2026 template (40 layouts). The three sub-brands resolve to their own
    embedded brand templates (deduplicated layout set + brand chrome).

    Search order:
      1. An uploaded .potx/.pptx on disk in the search dirs (CORPORATE only —
         a hand-supplied template always wins for the corporate deck).
      2. The split-module registry (takeda_template_<brand>.py, imported).
      3. Back-compat: a combined takeda_templates.py, if present.
      4. CORPORATE only: download Takeda_Slide_Template_EN.potx from GitHub
         via takeda_remote — so a bare Deck() works in both Mode A and Mode B.
    """
    brand = (brand or "CORPORATE").upper()

    # For the corporate deck, an uploaded template on disk takes priority.
    if brand == "CORPORATE":
        search_dirs = ["/mnt/project", "/home/claude", "/mnt/user-data/uploads"]
        stems = ["Takeda_Slide_Template_EN", "Takeda_Slide_Template",
                 "Takeda_Template", "Takeda_template"]
        exts = [".potx", ".pptx"]
        for d in search_dirs:
            for stem in stems:
                for ext in exts:
                    c = os.path.join(d, stem + ext)
                    if os.path.exists(c):
                        return c
        import glob as _glob
        for d in search_dirs:
            for ext in exts:
                hits = sorted(_glob.glob(os.path.join(d, "*akeda*" + ext)))
                if hits:
                    return hits[0]

    # All brands: decode from a loaded split template module via the shared registry.
    # Each takeda_template_<brand>.py registers itself into builtins on import.
    try:
        import builtins as _bi
        reg = getattr(_bi, "_TAKEDA_TEMPLATE_REGISTRY", {})
        if brand in reg:
            return _template_path_from_registry(brand)
    except Exception:
        pass

    # Back-compat: a single combined takeda_templates.py, if present.
    try:
        import takeda_templates as _tt
        return _tt.template_path(brand)
    except Exception:
        pass

    # Legacy /tmp cache from an earlier run, if present.
    _CACHE = "/tmp/takeda_template_cache.potx"
    if os.path.exists(_CACHE):
        return _CACHE

    # CORPORATE final fallback: download from GitHub via takeda_remote.
    # This makes a bare Deck() work in BOTH modes (Mode A authoring and the
    # Mode B clean canvas) without importing takeda_template_corporate first.
    if brand == "CORPORATE":
        try:
            import takeda_remote as _remote
            return _remote.get("Takeda_Slide_Template_EN.potx")
        except Exception:
            pass

    raise FileNotFoundError(
        f"Template for brand {brand!r} not found. Import the split module "
        f"takeda_template_{brand.lower()}, or place the template file in "
        f"/mnt/project/ (corporate: Takeda_Slide_Template_EN.potx), or check "
        f"that GitHub is reachable (takeda_remote.TEMPLATE_BASE_URL)."
    )

def _template_path_from_registry(brand):
    """Decode a brand template from the builtins registry populated by split modules."""
    import builtins as _bi, base64 as _b64, gzip as _gz
    reg = _bi._TAKEDA_TEMPLATE_REGISTRY
    cache_dir = "/tmp/takeda_templates"; os.makedirs(cache_dir, exist_ok=True)
    path = os.path.join(cache_dir, f"takeda_{brand.lower()}.pptx")
    if os.path.exists(path):
        return path
    val = reg[brand]
    if val == "__DISK__":
        import importlib, shutil
        mod = importlib.import_module(f"takeda_template_{brand.lower()}")
        shutil.copyfile(mod.template_path(brand), path)
        return path
    with open(path, "wb") as f:
        f.write(_gz.decompress(_b64.b64decode(val)))
    return path

# ── DECK CLASS ───────────────────────────────────────────────
class Deck:
    """Thin wrapper around python-pptx Presentation, pre-loaded with Takeda template."""

    def __init__(self, template_path=None, keep_template_slides=False, brand="CORPORATE"):
        self.brand = (brand or "CORPORATE").upper()
        path = template_path or _find_template(self.brand)
        path = Path(path)
        if path.suffix.lower() == ".potx":
            pptx_bytes = _potx_to_pptx_bytes(path)
            self.prs = Presentation(io.BytesIO(pptx_bytes))
        else:
            self.prs = Presentation(str(path))
        # 2026 template ships ~36 showcase slides (the prior template shipped 0).
        # Purge them so a generated deck starts empty. Layouts/masters are untouched.
        if not keep_template_slides:
            self._purge_slides()

    def _purge_slides(self):
        """Remove any pre-existing slides from the loaded template, leaving layouts intact."""
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            try:
                self.prs.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)

    # Ground-truth registry generated from the actual template (takeda_layouts.json):
    # semantic key -> {"master": mi, "layout": li, "name": str}. Resolves across ALL
    # masters by NAME, as the docs always promised. Falls back to the legacy index
    # registry (LAYOUTS) only if the JSON map is absent.
    _LAYOUT_MAP = None

    # Layout map is embedded here — no external JSON file needed.
    # Regenerate by running: python takeda.py --regen-layouts <template.potx>
    _LAYOUT_MAP_DATA = {
    "TITLE_PAGE_PHOTO": {
        "name": "Title Page 1",
        "master": 0,
        "layout": 0
    },
    "TITLE_PAGE_4": {
        "name": "4_Title Page 1",
        "master": 0,
        "layout": 1
    },
    "TITLE_PAGE_3": {
        "name": "3_Title Page 1",
        "master": 0,
        "layout": 2
    },
    "TITLE_PAGE_NO_PHOTO": {
        "name": "1_Title Page 1",
        "master": 0,
        "layout": 3
    },
    "DIVIDER_2": {
        "name": "2_Divider 1",
        "master": 0,
        "layout": 4
    },
    "DIVIDER_1": {
        "name": "1_Divider 1",
        "master": 0,
        "layout": 5
    },
    "DIVIDER_3": {
        "name": "3_Divider 1",
        "master": 0,
        "layout": 6
    },
    "ONE_COLUMN": {
        "name": "Standard 1-Column Text",
        "master": 0,
        "layout": 7
    },
    "ONE_COLUMN_WAVE": {
        "name": "1_Standard 1-Column Text",
        "master": 0,
        "layout": 8
    },
    "END_SLIDE": {
        "name": "End Slide",
        "master": 0,
        "layout": 9
    },
    "BLANK": {
        "name": "Blank slide",
        "master": 0,
        "layout": 10
    },
    "TWO_COLUMN_ADVANCED": {
        "name": "Advanced 2-Column Text",
        "master": 0,
        "layout": 11
    },
    "THREE_COLUMN": {
        "name": "Advanced 3-Column Text",
        "master": 0,
        "layout": 12
    },
    "FOUR_COLUMN": {
        "name": "Advanced 4-Column Text",
        "master": 0,
        "layout": 13
    },
    "IMAGE_TWO_THIRDS": {
        "name": "Advanced 2/3 Image",
        "master": 0,
        "layout": 14
    },
    "IMAGE_HALF": {
        "name": "Advanced _1/2 image",
        "master": 0,
        "layout": 15
    },
    "IMAGE_HALF_FULLBLEED": {
        "name": "Advanced 1/2 Image (Full Bleed)",
        "master": 0,
        "layout": 16
    },
    "IMAGE_ONE_THIRD": {
        "name": "Advanced 1/3 Image",
        "master": 0,
        "layout": 17
    },
    "IMAGE_ONE_THIRD_FULLBLEED": {
        "name": "Advanced 1/3 Image (Full Bleed)",
        "master": 0,
        "layout": 18
    },
    "IMAGE_FOUR_COLUMN": {
        "name": "Advanced Image 4-Column",
        "master": 0,
        "layout": 19
    },
    "CHART_FULL": {
        "name": "Advanced Chart Full Width",
        "master": 0,
        "layout": 20
    },
    "CHART_TWO_THIRDS": {
        "name": "Advanced Chart 2/3",
        "master": 0,
        "layout": 21
    },
    "CHART_TWO_COLUMN": {
        "name": "Advanced Chart 2 Column",
        "master": 0,
        "layout": 22
    },
    "BIG_PICTURE": {
        "name": "Big Picture",
        "master": 0,
        "layout": 23
    },
    "PATTERN_1_GRADIENT": {
        "name": "pattern",
        "master": 1,
        "layout": 0
    },
    "PATTERN_1_INVERTED": {
        "name": "pattern",
        "master": 1,
        "layout": 1
    },
    "PATTERN_1_SUBTLE_1": {
        "name": "pattern",
        "master": 1,
        "layout": 2
    },
    "PATTERN_1_SUBTLE_2": {
        "name": "pattern",
        "master": 1,
        "layout": 3
    },
    "PATTERN_2_GRADIENT": {
        "name": "pattern",
        "master": 2,
        "layout": 0
    },
    "PATTERN_2_INVERTED": {
        "name": "pattern",
        "master": 2,
        "layout": 1
    },
    "PATTERN_2_SUBTLE_1": {
        "name": "pattern",
        "master": 2,
        "layout": 2
    },
    "PATTERN_2_SUBTLE_2": {
        "name": "pattern",
        "master": 2,
        "layout": 3
    },
    "PATTERN_3_GRADIENT": {
        "name": "pattern",
        "master": 3,
        "layout": 0
    },
    "PATTERN_3_INVERTED": {
        "name": "pattern",
        "master": 3,
        "layout": 1
    },
    "PATTERN_3_SUBTLE_1": {
        "name": "pattern",
        "master": 3,
        "layout": 2
    },
    "PATTERN_3_SUBTLE_2": {
        "name": "pattern",
        "master": 3,
        "layout": 3
    },
    "PATTERN_4_GRADIENT": {
        "name": "pattern",
        "master": 4,
        "layout": 0
    },
    "PATTERN_4_INVERTED": {
        "name": "pattern",
        "master": 4,
        "layout": 1
    },
    "PATTERN_4_SUBTLE_1": {
        "name": "pattern",
        "master": 4,
        "layout": 2
    },
    "PATTERN_4_SUBTLE_2": {
        "name": "pattern",
        "master": 4,
        "layout": 3
    }
}

    @classmethod
    def _load_layout_map(cls):
        if cls._LAYOUT_MAP is None:
            cls._LAYOUT_MAP = cls._LAYOUT_MAP_DATA
        return cls._LAYOUT_MAP

    def _layout(self, key_or_index):
        if not isinstance(key_or_index, str):
            return self.prs.slide_layouts[int(key_or_index)]
        key = key_or_index
        if key in RETIRED_LAYOUTS:
            key = RETIRED_LAYOUTS[key]
        m = self._load_layout_map()
        if key in m:                       # ground-truth: cross-master by (master,layout)
            e = m[key]
            return self.prs.slide_masters[e["master"]].slide_layouts[e["layout"]]
        # Fallback: resolve the registry NAME by searching every master
        if key in LAYOUTS:
            want = LAYOUTS[key][1].strip()
            for master in self.prs.slide_masters:
                for L in master.slide_layouts:
                    if L.name.strip() == want:
                        return L
            return self.prs.slide_layouts[LAYOUTS[key][0]]  # last-resort index
        raise KeyError(f"Unknown layout key {key_or_index!r}")

    def pattern_layout(self, key):
        """Resolve a Pattern background layout (decorative full-bleed) by key.
        Pattern layouts live on masters 1-4, not master 0, so we index the
        master's layout list directly. Use for TITLE / DIVIDER slides only."""
        if key not in PATTERN_LAYOUTS:
            raise KeyError(f"Unknown pattern layout {key!r}")
        mi, li = PATTERN_LAYOUTS[key]
        return self.prs.slide_masters[mi].slide_layouts[li]

    def add_pattern_slide(self, key):
        """Add a slide on a Pattern background layout (title/divider use only)."""
        return self.prs.slides.add_slide(self.pattern_layout(key))

    def add_slide(self, layout_key):
        return self.prs.slides.add_slide(self._layout(layout_key))

    def add_slide_by_name(self, layout_name):
        """Add a slide by the layout's display NAME, searching all masters.
        This is how brand layouts (BRAND_LAYOUTS[brand]) are inserted — they
        don't have semantic keys in LAYOUTS. Case-insensitive; '&'/'&amp;' tolerant.
        Also applies the brand chrome's whitespace-contingent logo if applicable."""
        want = layout_name.replace("&amp;", "&").strip().lower()
        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name.replace("&amp;", "&").strip().lower() == want:
                    return self.prs.slides.add_slide(layout)
        raise KeyError(f"No layout named {layout_name!r} in the {self.brand} template. "
                       f"See BRAND_LAYOUTS.get('{self.brand}').")

    def list_layouts(self):
        """Return the layout names available in the loaded template (the 'new slide' set)."""
        return [layout.name for master in self.prs.slide_masters
                for layout in master.slide_layouts]

    @staticmethod
    def set_placeholder_text(slide, idx, text, preserve_formatting=True):
        ph = _find_placeholder(slide, idx)
        if ph is None or not ph.has_text_frame:
            return
        tf = ph.text_frame
        paragraphs = [text] if isinstance(text, str) else list(text)
        if not paragraphs:
            tf.text = ""; return
        if preserve_formatting and tf.paragraphs:
            first_p = tf.paragraphs[0]
            for r in list(first_p.runs):
                r._r.getparent().remove(r._r)
            run = first_p.add_run()
            run.text = paragraphs[0]
            base_p_xml = first_p._p
            for extra in list(tf.paragraphs[1:]):
                extra._p.getparent().remove(extra._p)
            for line in paragraphs[1:]:
                new_p = copy.deepcopy(base_p_xml)
                for child in list(new_p):
                    if child.tag.split("}")[-1] == "r":
                        new_p.remove(child)
                if first_p.runs:
                    new_run = copy.deepcopy(first_p.runs[0]._r)
                    t = new_run.find(qn("a:t"))
                    if t is not None: t.text = line
                    new_p.append(new_run)
                base_p_xml.addnext(new_p)
                base_p_xml = new_p
        else:
            tf.text = paragraphs[0]
            for line in paragraphs[1:]:
                p = tf.add_paragraph(); p.text = line

    # ── High-level slide helpers ─────────────────────────────
    def add_title(self, title, presenter="", role="", department="", date="", variant="no_photo"):
        # Title placeholder index differs by template version (legacy idx=0 CENTER_TITLE,
        # 2026 idx=27). Resolve by TYPE so it works on whichever template is loaded.
        key = {"no_photo":"TITLE_PAGE_NO_PHOTO","top_photo":"TITLE_PAGE_PHOTO",
               "photo":"TITLE_PAGE_PHOTO"}.get(variant, "TITLE_PAGE_NO_PHOTO")
        slide = self.add_slide(key)
        from pptx.enum.shapes import PP_PLACEHOLDER
        title_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                title_ph = ph; break
        if title_ph is not None:
            self.set_placeholder_text(slide, title_ph.placeholder_format.idx, title)
        else:                                    # last resort: try known indices
            self.set_placeholder_text(slide, 0, title)
            self.set_placeholder_text(slide, 27, title)
        for idx, val in ((26, presenter), (23, role), (24, department), (25, date)):
            if val:
                self.set_placeholder_text(slide, idx, val)
        return slide

    def add_section_title(self, title, subsection="", number=""):
        # 2026 template: divider layouts (1_Divider 1 default). Title idx=0,
        # subsection idx=27, number idx=28. Dividers now carry a pattern background.
        slide = self.add_slide("DIVIDER_1")
        self.set_placeholder_text(slide, 0, title)
        if subsection:
            self.set_placeholder_text(slide, 27, subsection)
        if number:
            self.set_placeholder_text(slide, 28, str(number))
        return slide

    def add_one_column(self, title, body, wave=False):
        slide = self.add_slide("ONE_COLUMN_WAVE" if wave else "ONE_COLUMN")
        self.set_placeholder_text(slide, 0, title)
        self.set_placeholder_text(slide, 15, body)

    def add_two_column(self, title, left, right, left_header="", right_header=""):
        slide = self.add_slide("TWO_COLUMN_ADVANCED")
        self.set_placeholder_text(slide, 0, title)
        if left_header or right_header:
            self.set_placeholder_text(slide, 14, left_header or right_header)
        self.set_placeholder_text(slide, 15, left)
        self.set_placeholder_text(slide, 16, right)

    def add_three_column(self, title, primary_title="", primary_text="", subtitle="", col1="", col2="", col3=""):
        slide = self.add_slide("THREE_COLUMN")
        self.set_placeholder_text(slide, 0, title)
        self.set_placeholder_text(slide, 21, primary_title); self.set_placeholder_text(slide, 20, primary_text)
        self.set_placeholder_text(slide, 14, subtitle)
        self.set_placeholder_text(slide, 15, col1); self.set_placeholder_text(slide, 18, col2); self.set_placeholder_text(slide, 19, col3)

    def add_four_column(self, title, primary_title="", primary_text="", subtitle="", col1="", col2="", col3="", col4=""):
        slide = self.add_slide("FOUR_COLUMN")
        self.set_placeholder_text(slide, 0, title)
        self.set_placeholder_text(slide, 22, primary_title); self.set_placeholder_text(slide, 21, primary_text)
        self.set_placeholder_text(slide, 14, subtitle)
        self.set_placeholder_text(slide, 15, col1); self.set_placeholder_text(slide, 18, col2)
        self.set_placeholder_text(slide, 19, col3); self.set_placeholder_text(slide, 20, col4)

    def add_image_slide(self, title, image_path=None, primary_title="", primary_text="", secondary_title="", secondary_text="", variant="two_thirds"):
        key = {"two_thirds":"IMAGE_TWO_THIRDS","half":"IMAGE_HALF","half_fullbleed":"IMAGE_HALF_FULLBLEED","one_third":"IMAGE_ONE_THIRD","one_third_fullbleed":"IMAGE_ONE_THIRD_FULLBLEED"}.get(variant,"IMAGE_TWO_THIRDS")
        slide = self.add_slide(key)
        self.set_placeholder_text(slide, 0, title)
        self.set_placeholder_text(slide, 21, primary_title); self.set_placeholder_text(slide, 20, primary_text)
        self.set_placeholder_text(slide, 14, secondary_title); self.set_placeholder_text(slide, 15, secondary_text)
        if image_path:
            pic_ph = _find_picture_placeholder(slide)
            if pic_ph and os.path.exists(image_path):
                pic_ph.insert_picture(image_path)

    def add_end_slide(self, text="Thank You"):
        slide = self.add_slide("END_SLIDE")
        self.set_placeholder_text(slide, 10, text)

    def save(self, path):
        path = Path(path); path.parent.mkdir(parents=True, exist_ok=True)
        uniquify_media_partnames(self.prs)   # prevent chrome/cloned-media partname collisions
        self.prs.save(str(path)); return path

# ── MODULE HELPERS ───────────────────────────────────────────
def _find_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None

def _find_picture_placeholder(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            return ph
    return None

# ===== Block 2: Mode B helper functions =====
# ── BODY PLACEHOLDER REMOVAL ─────────────────────────────────
# The Takeda legacy template defines these layout-inherited placeholders on
# ONE_COLUMN slides. If left empty they render as dashed "Click to add" boxes
# in PowerPoint. Remove them BEFORE cloning source content onto the slide.
#   idx=15  Body content placeholder
#   idx=16  Disclaimer / footer text placeholder
#   idx=3   Footer right placeholder
#   idx=4   Slide number placeholder
_PH_INDICES_TO_REMOVE = {15, 16, 3, 4}

def remove_body_placeholder(slide):
    """Remove all empty layout-inherited placeholders from a ONE_COLUMN slide."""
    spTree = slide.shapes._spTree
    removed = 0
    for sp in list(spTree.findall(qn("p:sp"))):
        ph = sp.find(".//" + qn("p:ph"))
        if ph is None:
            continue
        try:
            idx = int(ph.get("idx", "-1"))
        except ValueError:
            continue
        if idx in _PH_INDICES_TO_REMOVE:
            all_text = "".join(t.text or "" for t in sp.iter(qn("a:t"))).strip()
            if not all_text:
                spTree.remove(sp)
                removed += 1
    return removed > 0

# ── TITLE EXTRACTION ─────────────────────────────────────────
def extract_source_title(spTree):
    """Find source title text, remove the title shape, return text string."""
    for sp in list(spTree.findall(qn("p:sp"))):
        ph = sp.find(".//" + qn("p:ph"))
        if ph is not None and ph.get("type") in ("title", "ctrTitle"):
            text = "".join(t.text or "" for t in sp.iter(qn("a:t")))
            spTree.remove(sp)
            return text.strip()
    return None

# ── TITLE FONT SIZE ──────────────────────────────────────────
def pick_title_font_size(title_text, default_pt=28, min_pt=18):
    """Return largest font size ≤ default_pt that fits title in 2 lines.
    Char counts calibrated for 9.8" wide placeholder (constrained from 10.871")."""
    chars_per_line = {28: 50, 24: 59, 22: 64, 20: 70, 18: 78}
    for pt in sorted(chars_per_line.keys(), reverse=True):
        if pt > default_pt: continue
        if math.ceil(len(title_text) / chars_per_line[pt]) <= 2:
            return pt
    return min_pt

def set_title_font_size(slide, font_pt):
    """Override title placeholder (idx 0) font size at run + defRPr level."""
    for sh in slide.shapes:
        try:
            if sh.placeholder_format and sh.placeholder_format.idx == 0 and sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(font_pt)
                txBody = sh._element.find(qn("p:txBody"))
                if txBody is not None:
                    lstStyle = txBody.find(qn("a:lstStyle"))
                    if lstStyle is None: lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
                    lvl1 = lstStyle.find(qn("a:lvl1pPr"))
                    if lvl1 is None: lvl1 = etree.SubElement(lstStyle, qn("a:lvl1pPr"))
                    defRPr = lvl1.find(qn("a:defRPr"))
                    if defRPr is None: defRPr = etree.SubElement(lvl1, qn("a:defRPr"))
                    defRPr.set("sz", str(int(font_pt * 100)))
        except Exception:
            pass

# ── COLOR REMAPPING ──────────────────────────────────────────
def recolor_element(element, remap=None):
    """Remap all srgbClr values in an element tree to the Takeda palette.
    If remap dict not provided, auto-maps common source colors."""
    if remap is None:
        remap = {
            # ZS / common source colors → Takeda palette
            'EC7200': 'E1242A', 'F7941D': 'E1242A', 'FF6600': 'E1242A',  # oranges → red
            'FFA500': 'E1242A', 'FFD700': 'E1242A',
            '1A1628': '34373F', '333333': '34373F', '404040': '34373F',  # darks → charcoal
            '32A29B': 'A1A4AC', '009999': 'A1A4AC',  # teals → mid gray
            '0070C0': '34373F', '4472C4': '34373F', '5B9BD5': 'A1A4AC',  # blues → charcoal/gray
            'C00000': 'E1242A', 'FF0000': 'E1242A',  # reds → takeda red
            'D0D0D0': 'EDF2F4', 'D8D8D8': 'EDF2F4', 'E8E8E8': 'EDF2F4',  # light grays
            'F2F2F2': 'EDF2F4', 'EEEEEE': 'EDF2F4',
        }
    for el in element.iter(qn('a:srgbClr')):
        val = el.get('val', '').upper()
        if val in remap:
            el.set('val', remap[val])

def refont_element(element, target_font="Aptos"):
    """Replace all font references with the brand font (2026: Aptos)."""
    for tag in (qn('a:latin'), qn('a:ea'), qn('a:cs')):
        for el in element.iter(tag):
            el.set('typeface', target_font)
    for el in element.iter(qn('a:rPr')):
        for attr in ('fontFamily',):
            if el.get(attr):
                el.set(attr, target_font)

# ── FONT SIZE ENFORCEMENT ────────────────────────────────────
def enforce_min_font_size(spTree, min_pt=12, footnote_y_emu=6_200_000):
    """Bump every non-footnote body run below min_pt up to min_pt.

    ⚠️ DO NOT blanket-apply this when REFORMATTING an existing deck (Mode B).
    Source decks legitimately use small eyebrows (~11pt) and footnotes (~9-10pt);
    forcing them to 12pt makes text grow and OVERLAP the table/figure above it.
    When reformatting, PRESERVE original run sizes (only change font family + color).
    Use this only for NEW (Mode A) decks you author from scratch."""
    for sp in spTree.iter(qn("p:sp")):
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is not None:
                off = xfrm.find(qn("a:off"))
                if off is not None:
                    try:
                        if int(off.get("y", "0")) >= footnote_y_emu:
                            continue
                    except ValueError: pass
        for rPr in sp.iter(qn("a:rPr")):
            sz = rPr.get("sz")
            if sz:
                try:
                    if int(sz) < min_pt * 100:
                        rPr.set("sz", str(min_pt * 100))
                except ValueError: pass

# ── PARAGRAPH SPACING ────────────────────────────────────────
def add_paragraph_spacing(spTree, spacing_pt=600):
    """Add 6pt spcAft to multi-paragraph text frames (except last paragraph)."""
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for sp in spTree.iter(qn("p:sp")):
        txBody = sp.find(qn("p:txBody"))
        if txBody is None: continue
        paras = txBody.findall(qn("a:p"))
        if len(paras) <= 1: continue
        for p in paras[:-1]:
            pPr = p.find(qn("a:pPr"))
            if pPr is None:
                pPr = etree.SubElement(p, qn("a:pPr"))
                p.insert(0, pPr)
            # Remove spcBef if present
            for old in pPr.findall(qn("a:spcBef")):
                pPr.remove(old)
            # Set spcAft
            spcAft = pPr.find(qn("a:spcAft"))
            if spcAft is None:
                spcAft = etree.SubElement(pPr, qn("a:spcAft"))
            spcPts = spcAft.find(qn("a:spcPts"))
            if spcPts is None:
                spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
            cur_val = int(spcPts.get("val", "0"))
            if cur_val < spacing_pt:
                spcPts.set("val", str(spacing_pt))

# ── PRIVATE HELPERS + CONTRAST (v2: resolves theme/scheme colors & lum modifiers) ──
# Rewritten in Action 4. Resolves a:schemeClr and lumMod/lumOff/tint/shade so that
# tinted fills (e.g. accent1 + lumMod20/lumOff80 = pale pink) are read by their
# RENDERED color, not their dark base. Defaults inherited text per the gold-deck
# rule: dark/colored fill -> white text; light fill -> charcoal. Validated against
# 75 real Takeda slides: 0 harmful changes; repairs 100% of broken dark-on-dark.

def _hex_to_rgb(hex6):
    return (int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))

def _luminance(r, g, b):
    """Relative luminance 0..1 from 0-255 ints (perceptual weights)."""
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0

# Default theme color map (Takeda 2021 theme). schemeClr names → hex.
# Used when a fill/text references a theme slot instead of an explicit hex.
_THEME = {
    'tx1': '000000', 'dk1': '000000', 'bg1': 'FFFFFF', 'lt1': 'FFFFFF',
    'tx2': '34373F', 'dk2': '34373F', 'bg2': 'EDF2F4', 'lt2': 'EDF2F4',
    'accent1': 'E1242A', 'accent2': '34373F', 'accent3': '7D8394',
    'accent4': '2E424B', 'accent5': 'A1A4AC', 'accent6': '891515',
    'hlink': '0070C0', 'folHlink': '954F72',
}

def _apply_lum_mods(hex6, clr_el):
    """Apply a:lumMod / a:lumOff / a:tint / a:shade / a:alpha modifiers to a base
    hex color, returning the effective rendered hex. This is what turns e.g.
    accent1 (red) + lumMod20%/lumOff80% into the pale pink actually displayed.
    Without this, a tinted fill is mistaken for its dark base color."""
    r, g, b = _hex_to_rgb(hex6)
    def _get(tag):
        e = clr_el.find(qn(tag))
        return int(e.get('val')) / 100000.0 if e is not None and e.get('val') else None
    lumMod = _get('a:lumMod')
    lumOff = _get('a:lumOff')
    tint   = _get('a:tint')
    shade  = _get('a:shade')
    if lumMod is not None:
        r, g, b = r * lumMod, g * lumMod, b * lumMod
    if lumOff is not None:
        r = r + 255 * lumOff
        g = g + 255 * lumOff
        b = b + 255 * lumOff
    if tint is not None:   # tint moves toward white
        r = r * tint + 255 * (1 - tint)
        g = g * tint + 255 * (1 - tint)
        b = b * tint + 255 * (1 - tint)
    if shade is not None:  # shade moves toward black
        r, g, b = r * shade, g * shade, b * shade
    clamp = lambda v: max(0, min(255, int(round(v))))
    return f"{clamp(r):02X}{clamp(g):02X}{clamp(b):02X}"

def _resolve_srgb_or_scheme(el):
    """Given an element that may contain a:srgbClr or a:schemeClr, return the
    EFFECTIVE hex (after luminance/tint modifiers), or None if unresolvable."""
    if el is None:
        return None
    srgb = el.find(qn('a:srgbClr'))
    if srgb is not None:
        base = srgb.get('val', '000000')
        return _apply_lum_mods(base, srgb)
    scheme = el.find(qn('a:schemeClr'))
    if scheme is not None:
        name = scheme.get('val', '')
        # phClr means "inherit placeholder color" — not resolvable here
        if name in _THEME:
            return _apply_lum_mods(_THEME[name], scheme)
    return None

def _fill_hex(shape):
    """Resolve a shape's solid fill to a hex string, handling srgbClr AND schemeClr.
    Scans ONLY shape-level spPr (never a:rPr — that is run text color, not fill).
    Returns hex or None if no solid fill found."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            hx = _resolve_srgb_or_scheme(sf)
            if hx:
                return hx
    return None

def _run_hex(run):
    """Resolve a run's font color to hex, handling explicit rgb, srgbClr, schemeClr.
    Returns hex or None if truly unresolvable (inherited from placeholder)."""
    try:
        rgb = run.font.color.rgb
        if rgb is not None:
            return str(rgb)
    except (AttributeError, TypeError):
        pass
    rPr = run._r.find(qn('a:rPr'))
    if rPr is not None:
        sf = rPr.find(qn('a:solidFill'))
        if sf is not None:
            hx = _resolve_srgb_or_scheme(sf)
            if hx:
                return hx
    return None

def _lum_of_hex(hex6):
    return _luminance(*_hex_to_rgb(hex6))

def _contrast_ratio(hex_a, hex_b):
    """WCAG contrast ratio between two hex colors (1.0 .. 21.0)."""
    def _lin(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    def _rel_lum(hex6):
        r, g, b = _hex_to_rgb(hex6)
        return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)
    la, lb = _rel_lum(hex_a), _rel_lum(hex_b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)

# ── BACKWARD-COMPAT shims (other code calls these) ──────────
def _get_fill_luminance(shape):
    hx = _fill_hex(shape)
    return _lum_of_hex(hx) if hx else None

def _get_run_luminance(run):
    hx = _run_hex(run)
    return _lum_of_hex(hx) if hx else None

def _has_dark_fill(shape):
    lum = _get_fill_luminance(shape)
    return lum is not None and lum < 0.45

# ── ENFORCE CONTRAST (v2) ───────────────────────────────────
WHITE = 'FFFFFF'
CHARCOAL = '34373F'
_MIN_RATIO = 3.0          # below this, text is hard to read → fix it
_SLIDE_BG = 'FFFFFF'      # assumed slide background when fill unresolved

def _fix_runs(paragraphs, fill_hex):
    """Ensure readable text against fill_hex (None → slide background, white).

    Gold-deck rule: a dark/colored fill takes WHITE text; a light fill takes
    CHARCOAL text. We compute the rule's target, then only change a run when its
    current (resolved) color is unreadable — i.e. contrast < _MIN_RATIO against
    the fill. Inherited/unresolved colors are always set to the rule target."""
    bg = fill_hex or _SLIDE_BG
    bg_is_dark = _lum_of_hex(bg) < 0.5
    rule_target = WHITE if bg_is_dark else CHARCOAL
    for para in paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            txt_hex = _run_hex(run)
            if txt_hex is None:
                # inherited/theme color we can't read → apply the rule target
                run.font.color.rgb = RGBColor.from_string(rule_target)
                continue
            if bg_is_dark:
                # Gold-deck rule: dark/colored fill ALWAYS takes white text.
                # Anything not already white is wrong (e.g. black-on-red).
                if txt_hex.upper() != WHITE:
                    run.font.color.rgb = RGBColor.from_string(WHITE)
            else:
                # Light fill: charcoal is canonical, but an on-brand accent
                # (red/crimson) that still clears the contrast floor is allowed.
                if _contrast_ratio(txt_hex, bg) < _MIN_RATIO:
                    run.font.color.rgb = RGBColor.from_string(CHARCOAL)

def enforce_contrast(slide):
    """Resolve inherited/theme colors and guarantee readable text on every fill.
    Gold-deck rule: dark/colored fill → white text; light fill → charcoal text.
    Uses WCAG contrast ratio (<3.0 triggers a fix) and defaults unresolved text
    against the fill (or slide background when the fill is also unresolved)."""
    for shape in slide.shapes:
        fill_hex = _fill_hex(shape)
        if shape.has_text_frame:
            # A text box with no fill still needs checking against the slide bg
            _fix_runs(shape.text_frame.paragraphs, fill_hex)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_fill = _resolve_srgb_or_scheme(
                        cell._tc.find(qn('a:tcPr')).find(qn('a:solidFill'))
                        if cell._tc.find(qn('a:tcPr')) is not None
                        and cell._tc.find(qn('a:tcPr')).find(qn('a:solidFill')) is not None
                        else None
                    )
                    _fix_runs(cell.text_frame.paragraphs, cell_fill)

def _is_takeda_placeholder(shape):
    """True if shape is a layout-defined placeholder (part of Takeda template chrome)."""
    try:
        return shape.placeholder_format is not None
    except Exception:
        return False

# ── NORMALIZE TITLE TEXT ────────────────────────────────────
def normalize_title_text(text):
    """Collapse whitespace, strip pagination suffixes, trim."""
    import re
    text = re.sub(r'[\x0b\t\r\n]+', ' ', text)
    text = re.sub(r'\s{2,}', ' ', text)
    text = re.sub(r'\s*Page\s+\d+\s*(of|/)\s*\d+\s*', '', text, flags=re.I)
    text = re.sub(r'\s*Slide\s+\d+\s*', '', text, flags=re.I)
    return text.strip()

# ── STRIP SOURCE CHROME ─────────────────────────────────────
def strip_source_chrome(slide):
    """Remove source header bars, pagination, logos, footers that conflict with Takeda chrome."""
    import re
    shapes_to_remove = []
    for shape in slide.shapes:
        # Guard: cloned shapes can inherit geometry (None). Skip those safely.
        if None in (shape.left, shape.top, shape.width, shape.height):
            continue
        l = shape.left / 914400
        t = shape.top / 914400
        w = shape.width / 914400
        h = shape.height / 914400
        # Full-width dark bars in title zone
        if t < 1.10 and w > 10.0 and h < 0.80:
            if _has_dark_fill(shape):
                shapes_to_remove.append(shape)
                continue
        # Shapes in Logo Exclusion Zone
        if l > 11.0 and t < 1.10:
            shapes_to_remove.append(shape)
            continue
        # Source pagination text
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if re.match(r'^(Page\s+\d+\s*(of|/)\s*\d+|Slide\s+\d+|\d{1,3}\s*$)', txt, re.I):
                shapes_to_remove.append(shape)
                continue
        # Source footer zone shapes
        if t > 7.0 and not _is_takeda_placeholder(shape):
            shapes_to_remove.append(shape)
            continue
    spTree = slide.shapes._spTree
    for shape in shapes_to_remove:
        sp = shape._element
        if sp.getparent() is not None:
            sp.getparent().remove(sp)

# ── ENFORCE CONTENT ZONES ───────────────────────────────────
def enforce_content_zones(slide):
    """Shift shapes that intrude into the left-accent or logo exclusion zones.
    NEVER changes a shape's width/height (that distorts images and squashes
    shapes) — only repositions. For top-right logo-zone intrusions, the shape
    is nudged DOWN below the logo band rather than narrowed."""
    MIN_LEFT = Inches(0.663)
    for shape in slide.shapes:
        if _is_takeda_placeholder(shape):
            continue
        l = shape.left; t = shape.top; w = shape.width
        if l is None or t is None or w is None:
            continue
        t_in = t / 914400
        # Left-accent zone: shift right (no resize)
        if t_in >= 1.10 and l < MIN_LEFT:
            shape.left = MIN_LEFT
        # Logo zone (top-right): if a non-placeholder shape sits under the logo,
        # move it DOWN below the logo band — do NOT shrink its width.
        if t_in < 1.10:
            right_edge = (l + w) / 914400
            if right_edge > 11.0:
                shape.top = Inches(1.20)   # drop below logo/title band

def enforce_icon_clearance(slide, pad_in=0.06):
    """Prevent icons from overlapping body text.

    An 'icon' here = a small (<0.8") picture / freeform / group that is NOT a
    placeholder and NOT the blade/cover chrome. For each icon that overlaps a
    text-bearing shape, move the icon to the text block's left margin (just
    outside it) if room exists, otherwise nudge it above the text. This fixes
    the 'label text hidden behind icon' bug (e.g. card icons landing on the
    label). Idempotent enough to run once at the end of a slide build."""
    def _box(sh):
        return ((sh.left or 0)/914400, (sh.top or 0)/914400,
                (sh.width or 0)/914400, (sh.height or 0)/914400)
    def _overlap(a, b):
        ax, ay, aw, ah = a; bx, by, bw, bh = b
        return not (ax+aw <= bx or bx+bw <= ax or ay+ah <= by or by+bh <= ay)

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    text_shapes = []
    icons = []
    for sh in slide.shapes:
        nm = (sh.name or '')
        if nm in ('TakedaBlade', 'StripCover'):
            continue
        has_text = bool(getattr(sh, 'has_text_frame', False) and sh.text_frame.text.strip())
        bx, by, bw, bh = _box(sh)
        is_iconish = (sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.FREEFORM,
                                        MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE)
                      and bw <= 0.8 and bh <= 0.8 and not has_text)
        if is_iconish:
            icons.append(sh)
        elif has_text:
            text_shapes.append(sh)

    moved = 0
    for icon in icons:
        ib = _box(icon)
        for txt in text_shapes:
            tb = _box(txt)
            if not _overlap(ib, tb):
                continue
            ix, iy, iw, ih = ib; tx, ty, tw, th = tb
            # Preferred: place icon just left of the text block, vertically centered to it
            new_left_in = tx - iw - pad_in
            if new_left_in >= 0.663:
                icon.left = Inches(new_left_in)
                icon.top  = Inches(ty + (th - ih) / 2.0)
            else:
                # No room at left → place icon just above the text block
                icon.top = Inches(max(1.20, ty - ih - pad_in))
            moved += 1
            break
    return moved

# ── PALETTE AUDIT ────────────────────────────────────────────

# ── BLADE ACCENT SHAPE (2026 — replaces layout rectangle) ───
# The legacy template's "Rectangle 20" (L=0.34" T=0.18" W=0.059" H=0.75",
# filled Takeda Red) lives on the layout layer and cannot be removed per-slide
# via python-pptx. The 2026 refresh replaced it with a curved Bezier "blade"
# (quarter-leaf) at L=0.34" T=0.34" W=0.157" H=0.477".
#
# Fix: on every body slide, insert (1) a white cover rectangle oversized by
# 0.04" on each side over the layout strip to fully mask it, then (2) the
# blade shape on top. Both shapes are locked/non-selectable so they behave
# like layout chrome.

def _make_blade_xml(shape_id):
    """Build 2026 Takeda Bezier blade shape XML (quarter-leaf curve)."""
    PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    BX = int(0.34 * 914400); BY = int(0.34 * 914400)
    BW = int(0.157 * 914400); BH = int(0.477 * 914400)
    return f'''<p:sp xmlns:p="{PML}" xmlns:a="{DML}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="TakedaBlade"/>
    <p:cNvSpPr><a:spLocks noGrp="1" noSelect="1" noMove="1" noResize="1" noDrilldown="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{BX}" y="{BY}"/><a:ext cx="{BW}" cy="{BH}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="{BW}" b="{BH}"/>
      <a:pathLst>
        <a:path w="{BW}" h="{BH}">
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:cubicBezTo>
            <a:pt x="{int(BW*0.4)}" y="{int(BH*0.05)}"/>
            <a:pt x="{BW}" y="{int(BH*0.35)}"/>
            <a:pt x="{BW}" y="{int(BH*0.5)}"/>
          </a:cubicBezTo>
          <a:cubicBezTo>
            <a:pt x="{BW}" y="{int(BH*0.65)}"/>
            <a:pt x="{int(BW*0.4)}" y="{int(BH*0.95)}"/>
            <a:pt x="0" y="{BH}"/>
          </a:cubicBezTo>
          <a:close/>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:solidFill><a:srgbClr val="E1242A"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:style>
    <a:lnRef idx="0"><a:schemeClr val="accent1"/></a:lnRef>
    <a:fillRef idx="0"><a:schemeClr val="accent1"/></a:fillRef>
    <a:effectRef idx="0"><a:schemeClr val="accent1"/></a:effectRef>
    <a:fontRef idx="minor"><a:schemeClr val="tx1"/></a:fontRef>
  </p:style>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''

def _make_cover_rect_xml(shape_id, L_in, T_in, W_in, H_in):
    """White rectangle sized to fully mask a detected red strip/bar.
    Caller passes the ACTUAL geometry of the strip to cover (plus margin),
    so this works regardless of strip height (1-line vs 2-line title)."""
    PML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    CL = int(L_in * 914400); CT = int(T_in * 914400)
    CW = int(W_in * 914400); CH = int(H_in * 914400)
    return f'''<p:sp xmlns:p="{PML}" xmlns:a="{DML}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="StripCover"/>
    <p:cNvSpPr><a:spLocks noGrp="1" noSelect="1" noMove="1" noResize="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{CL}" y="{CT}"/><a:ext cx="{CW}" cy="{CH}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody>
</p:sp>'''

def _is_red_strip_shape(sh):
    """True if a shape is a thin tall red vertical bar in the title/left zone —
    i.e. legacy Takeda red accent strip (any height). Detects both the old
    layout rectangle and red bars cloned in from a Mode B source slide."""
    try:
        L = (sh.left or 0) / 914400; W = (sh.width or 0) / 914400
        H = (sh.height or 0) / 914400; Tp = (sh.top or 0) / 914400
    except Exception:
        return False
    if not (W < 0.18 and H > 0.35 and L < 1.2 and Tp < 1.6):
        return False  # not a thin tall bar near the top-left
    # Must be red (E1242A family) via srgb or accent1 scheme
    reds = {'E1242A', 'C00000', 'FF0000', '891515'}
    for e in sh._element.iter(qn('a:srgbClr')):
        if e.get('val', '').upper() in reds:
            return True
    for e in sh._element.iter(qn('a:schemeClr')):
        if e.get('val') == 'accent1':
            return True
    return False

def _slide_has_blade(slide):
    """True if the 2026 blade (layout or slide) is already present."""
    # Check slide-level shapes we may have added
    for sp in slide.shapes._spTree.findall(qn('p:sp')):
        nvPr = sp.find('.//' + qn('p:cNvPr'))
        if nvPr is not None and nvPr.get('name', '') == 'TakedaBlade':
            return True
    # Check the layout: the 2026 content layouts ship the blade as a custGeom
    # accent1 freeform near the top-left. If present, we must NOT add another.
    try:
        for sh in slide.slide_layout.shapes:
            if sh.is_placeholder:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                L = (sh.left or 0) / 914400; W = (sh.width or 0) / 914400
                if L < 0.6 and W < 0.4:
                    return True
    except Exception:
        pass
    return False

def add_blade_accent(slide, base_id=800):
    """Ensure the slide shows the 2026 Bezier blade accent and NO straight red strip.

    Handles three cases robustly:
      (a) 2026 template, blade already on layout  → do nothing (no double blade).
      (b) Legacy template / Mode B source carrying a straight red strip
          (any height) → cover it with a white rect sized to the actual strip,
          then draw the blade.
      (c) Slide with neither → just draw the blade.
    Idempotent. Call after constrain_title_placeholder(), before recolor_element()."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    spTree = slide.shapes._spTree

    # Idempotency: if we already added a blade/cover to this slide, stop.
    for sp in spTree.findall(qn('p:sp')):
        nvPr = sp.find('.//' + qn('p:cNvPr'))
        if nvPr is not None and nvPr.get('name', '') in ('TakedaBlade', 'StripCover'):
            return

    # 1) Remove any straight red strip that lives ON THE SLIDE (e.g. cloned from a
    #    Mode B source). Direct removal is cleaner than covering when we own the shape.
    for sh in list(slide.shapes):
        if not sh.is_placeholder and _is_red_strip_shape(sh):
            sh._element.getparent().remove(sh._element)

    # 2) If a red strip lives on the LAYOUT (legacy template), cover it precisely.
    #    Size the cover to the strip's real geometry + margin (handles 2-line titles).
    covered = False
    try:
        for sh in slide.slide_layout.shapes:
            if sh.is_placeholder:
                continue
            if _is_red_strip_shape(sh):
                L = (sh.left or 0) / 914400; Tp = (sh.top or 0) / 914400
                W = (sh.width or 0) / 914400; H = (sh.height or 0) / 914400
                cover = etree.fromstring(_make_cover_rect_xml(
                    base_id, L - 0.04, Tp - 0.04, W + 0.10, H + 0.08))
                spTree.insert(1, cover)
                covered = True
    except Exception:
        pass

    # 3) If the 2026 blade is already present on the layout AND we did not have to
    #    cover a legacy strip, do nothing — the layout blade is correct as-is.
    if _slide_has_blade(slide) and not covered:
        return

    # 4) Draw the blade (covered legacy strip, or layout had no blade).
    blade = etree.fromstring(_make_blade_xml(base_id + 1))
    spTree.insert(2, blade)

# ── PALETTE AUDIT ────────────────────────────────────────────
def audit_colors(prs):
    """Return list of (slide_num, hex_value) for any off-palette srgbClr."""
    off = []
    for si, slide in enumerate(prs.slides):
        for el in slide._element.iter(qn('a:srgbClr')):
            val = el.get('val', '').upper()
            if val and val not in TAKEDA_PALETTE:
                off.append((si + 1, val))
    return off

def remap_colors(prs, remap_dict):
    """Remap off-palette colors in-place across all slides."""
    count = 0
    for slide in prs.slides:
        for el in slide._element.iter(qn('a:srgbClr')):
            val = el.get('val', '').upper()
            if val in remap_dict:
                el.set('val', remap_dict[val])
                count += 1
    return count

# ── CHART RECOLORING ─────────────────────────────────────────
def recolor_charts(prs, remap=None):
    """Recolor chart series in ppt/charts/*.xml files.
    Charts store colors separately from the slide spTree."""
    if remap is None:
        remap = {
            'EC7200': 'E1242A', 'F7941D': 'E1242A', '32A29B': 'A1A4AC',
            '1A1628': '34373F', '0070C0': '34373F', '4472C4': '34373F',
            '5B9BD5': 'A1A4AC', 'C00000': 'E1242A',
        }
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                chart_part = shape.chart.part
                chart_xml = chart_part._element
                for el in chart_xml.iter(qn('a:srgbClr')):
                    val = el.get('val', '').upper()
                    if val in remap:
                        el.set('val', remap[val])
                        count += 1
    return count

# ── TITLE PLACEHOLDER CONSTRAINT ────────────────────────────
def constrain_title_placeholder(slide):
    """Shrink title placeholder (idx=0) width to prevent logo overlap,
    and align to the 2026 template title position.

    IMPORTANT: always set BOTH position AND size. A title that gets an <a:off>
    (position) but no <a:ext> (size) is read by PowerPoint as zero width, which
    renders the title one character per line (vertical stacking). LibreOffice masks
    this by falling back to the layout width, so it only shows in PowerPoint. Setting
    left/top/width/height together (below) writes a complete xfrm and avoids it.
    """
    for sh in slide.shapes:
        try:
            if sh.placeholder_format and sh.placeholder_format.idx == 0:
                sh.left  = Inches(0.751)   # 2026 template title left (was 0.663)
                sh.width = Inches(9.8)     # right edge clears the top-right logo
                sh.top   = Inches(0.210)   # 2026 template title top (was 0.177)
                sh.height = Inches(0.75)
                # Safety check: confirm a non-empty width landed (catches the zero-width bug)
                if sh.width is None or int(sh.width) <= 0:
                    sh.width = Inches(9.8)
                return
        except Exception:
            pass

# ===== Block 3: Relationship-aware shape cloning =====
def clone_shapes_with_rels(src_slide, new_slide):
    """Clone all shapes from source slide to new slide, properly copying
    image/chart/hyperlink relationships so PowerPoint needs no repair.
    Automatically strips think-cell OLE objects (cause repair warnings)."""
    
    src_spTree = src_slide.shapes._spTree
    new_spTree = new_slide.shapes._spTree
    cloneable = ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp')
    count = 0
    
    for src_el in list(src_spTree):
        tag = src_el.tag.split('}')[-1] if '}' in src_el.tag else src_el.tag
        if tag not in cloneable:
            continue
        # Skip think-cell / OLE remnants that cause PowerPoint "repair and removed" warnings.
        # TWO cases, both must be caught (a missed case discards the WHOLE slide on open):
        #   (a) a graphicFrame containing a think-cell <p:oleObj> (progId TCLayout / think-cell)
        #   (b) a graphicFrame whose <a:graphicData uri=".../ole"> is EMPTY (no child object) —
        #       the orphaned sibling left behind when a think-cell AlternateContent block is
        #       removed. PowerPoint cannot read an empty OLE frame and discards the slide spTree.
        if tag == 'graphicFrame' and _is_dead_ole_frame(src_el):
            continue
        cloned = copy.deepcopy(src_el)
        _fix_rels_in_element(cloned, src_slide, new_slide)
        _migrate_table_styles(cloned, src_slide, new_slide)
        new_spTree.append(cloned)
        count += 1
    return count

def _get_tablestyles_part(package):
    """Return the package-level /ppt/tableStyles.xml part, or None."""
    for part in package.iter_parts():
        if str(part.partname) == '/ppt/tableStyles.xml':
            return part
    return None

def _migrate_table_styles(element, src_slide, new_slide):
    """Ensure every <a:tableStyleId> a cloned table references exists in the
    TARGET package's tableStyles.xml. Table styles live in one package-level
    part, not on the shape, so cloning a table copies its style-GUID reference
    but not the style itself. A dangling GUID makes PowerPoint report the file
    as corrupt ('found a problem with content') and offer to repair — even
    though the XSD and LibreOffice both accept it. Copy any missing definition
    from the SOURCE package; GUIDs absent from source are PowerPoint built-ins
    and need no entry."""
    from lxml import etree as _etree
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    aq = '{%s}' % A
    ids = [e.text.strip() for e in element.iter(aq + 'tableStyleId')
           if e.text and e.text.strip()]
    if not ids:
        return
    dst_part = _get_tablestyles_part(new_slide.part.package)
    src_part = _get_tablestyles_part(src_slide.part.package)
    if dst_part is None or src_part is None:
        return
    dst_root = _etree.fromstring(dst_part.blob)
    src_root = _etree.fromstring(src_part.blob)
    present = {e.get('styleId') for e in dst_root.iter(aq + 'tblStyle')}
    src_map = {e.get('styleId'): e for e in src_root.iter(aq + 'tblStyle')}
    changed = False
    for g in ids:
        if g not in present and g in src_map:
            dst_root.append(_etree.fromstring(_etree.tostring(src_map[g])))
            present.add(g)
            changed = True
    if changed:
        dst_part._blob = _etree.tostring(dst_root, xml_declaration=True,
                                         encoding='UTF-8', standalone=True)

def _is_dead_ole_frame(gf):
    """True if a <p:graphicFrame> is a think-cell OLE object OR an orphaned empty-OLE
    frame — either of which triggers a PowerPoint repair-and-remove on open."""
    OLE_URI = 'http://schemas.openxmlformats.org/presentationml/2006/ole'
    # (a) think-cell oleObj anywhere inside
    for ole in gf.iter(qn('p:oleObj')):
        prog = ole.get('progId', '') or ''
        if 'TCLayout' in prog or 'think-cell' in prog.lower():
            return True
    # (b) empty OLE graphicData (uri says OLE but no child element)
    for gd in gf.iter(qn('a:graphicData')):
        if 'ole' in (gd.get('uri', '') or '').lower() and len(list(gd)) == 0:
            return True
    return False

def _fix_rels_in_element(element, src_slide, new_slide):
    """Walk all rId references in an element and re-map to new slide relationships."""
    RNS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    CHART_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    
    # 1. Image blips: <a:blip r:embed="rIdN"/>
    for blip in element.iter(qn('a:blip')):
        for attr in (qn('r:embed'), qn('r:link')):
            old_rId = blip.get(attr)
            if old_rId:
                new_rId = _copy_rel(old_rId, src_slide, new_slide)
                if new_rId: blip.set(attr, new_rId)
                else: blip.attrib.pop(attr, None)
    
    # 1b. SVG blips: <asvg:svgBlip r:embed="rIdN"/> nested in a blip's extLst.
    #     clone copies the vector layer but this rId is separate from the raster
    #     <a:blip> above; if it isn't remapped the reference dangles and PowerPoint
    #     reports the deck as corrupt. On failure, drop the whole SVG <a:ext> so the
    #     raster fallback still renders (never leave an svgBlip without a valid rId).
    SVG_NS = 'http://schemas.microsoft.com/office/drawing/2016/SVG/main'
    for svg in list(element.iter(f'{{{SVG_NS}}}svgBlip')):
        old_rId = svg.get(qn('r:embed')) or svg.get(qn('r:link'))
        new_rId = _copy_rel(old_rId, src_slide, new_slide) if old_rId else None
        if new_rId:
            if svg.get(qn('r:embed')): svg.set(qn('r:embed'), new_rId)
            else: svg.set(qn('r:link'), new_rId)
        else:
            ext = svg.getparent()
            while ext is not None and ext.tag != qn('a:ext'):
                ext = ext.getparent()
            if ext is not None and ext.getparent() is not None:
                ext.getparent().remove(ext)

    # 2. Charts: <c:chart r:id="rIdN"/>
    for chart_ref in element.iter(f'{{{CHART_NS}}}chart'):
        old_rId = chart_ref.get(qn('r:id'))
        if old_rId:
            new_rId = _copy_rel(old_rId, src_slide, new_slide)
            if new_rId: chart_ref.set(qn('r:id'), new_rId)
    
    # 3. Hyperlinks: <a:hlinkClick r:id="rIdN"/>
    for hlink in element.iter(qn('a:hlinkClick')):
        old_rId = hlink.get(qn('r:id'))
        if old_rId and old_rId.strip():
            new_rId = _copy_rel(old_rId, src_slide, new_slide)
            if new_rId: hlink.set(qn('r:id'), new_rId)
            else: hlink.attrib.pop(qn('r:id'), None)
    
    # 4. OLE objects
    for ole in element.iter(qn('p:oleObj')):
        old_rId = ole.get(qn('r:id'))
        if old_rId:
            new_rId = _copy_rel(old_rId, src_slide, new_slide)
            if new_rId: ole.set(qn('r:id'), new_rId)

def _copy_rel(old_rId, src_slide, new_slide):
    """Copy a single relationship from source to new slide. Returns new rId or None."""
    try:
        src_rel = src_slide.part.rels[old_rId]
    except KeyError:
        return None
    try:
        if src_rel.is_external:
            return new_slide.part.rels.get_or_add_ext_rel(src_rel.reltype, src_rel.target_ref)
        else:
            return new_slide.part.relate_to(src_rel.target_part, src_rel.reltype)
    except Exception:
        return None

# ===== Block 4: Post-save ZIP cleanup =====
def clean_pptx_zip(path):
    """Remove duplicate ZIP entries that cause repair warnings."""
    import zipfile, io
    with zipfile.ZipFile(path, "r") as zin:
        seen = set()
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename not in seen:
                    seen.add(item.filename)
                    zout.writestr(item, zin.read(item.filename))
        with open(path, "wb") as f:
            f.write(buf.getvalue())


def uniquify_media_partnames(prs):
    """Give every media part a UNIQUE partname so no two distinct images collide.

    Root cause of the "banner renders as a stray map" bug: `clone_shapes_with_rels`
    relates foreign source image parts into the deck, and those parts KEEP their source
    partnames (e.g. /ppt/media/image28.png). Chrome images added via `add_picture` are
    auto-named by python-pptx and can land on the SAME name. On save, python-pptx writes
    two ZIP entries with that name; `clean_pptx_zip` then keeps only one blob, so one of
    the two pictures (typically the FRUZAQLA banner) silently points at the wrong image.

    python-pptx rels reference the target PART OBJECT and compute the serialized target
    ref from `part.partname` at save time, so simply reassigning a colliding part's
    `.partname` transparently repairs every relationship that points to it. Must run
    BEFORE save (the losing blob is destroyed during serialization otherwise).
    """
    from pptx.opc.packuri import PackURI
    from collections import defaultdict
    pkg = prs.part.package
    parts = list(pkg.iter_parts())
    used = {str(p.partname) for p in parts}
    by_name = defaultdict(list)
    for p in parts:
        pn = str(p.partname)
        if pn.startswith("/ppt/media/"):
            by_name[pn].append(p)
    n = 1
    for pn, plist in by_name.items():
        if len(plist) <= 1:
            continue  # no collision
        ext = pn.rsplit(".", 1)[-1] if "." in pn else "png"
        for p in plist[1:]:                     # keep first, rename the rest
            while True:
                cand = f"/ppt/media/uniq{n}.{ext}"
                n += 1
                if cand not in used:
                    used.add(cand)
                    break
            p.partname = PackURI(cand)
    return prs


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — MODE B / BRAND CONVERSION HELPERS
# (Functions take brand colors as arguments — no per-brand constants here)
# ═══════════════════════════════════════════════════════════════════════════════

# --- from brand block [0] ---
def set_brand_title(slide, deck, title_text, brand_color_hex, brand_font="Aptos"):
    """Set the title as brand-colored text in a no-fill placeholder.
    NEVER draws a filled bar. Color is set explicitly so it can't render black."""
    deck.set_placeholder_text(slide, 0, title_text)
    for sh in slide.shapes:
        try:
            if sh.placeholder_format and sh.placeholder_format.idx == 0:
                # Guarantee NO fill on the title shape
                spPr = sh._element.find('.//' + qn('p:spPr'))
                if spPr is not None:
                    for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
                        existing = spPr.find(qn(tag))
                        if existing is not None:
                            spPr.remove(existing)
                    if spPr.find(qn('a:noFill')) is None:
                        spPr.append(spPr.makeelement(qn('a:noFill'), {}))
                # Set text color EXPLICITLY (prevents black inheritance)
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor.from_string(brand_color_hex)
                        run.font.name = brand_font
                        run.font.bold = True
                        run.font.size = Pt(20)
        except Exception:
            pass
    constrain_title_placeholder(slide)

# --- from brand block [2] ---
def neutralize_takeda_layout_chrome(slide):
    """Remove Takeda corporate furniture (logo image + red accent strip) injected by the
    layout/master, so the brand chrome (e.g. FRUZAQLA header strip) is the only branding.
    Idempotent; safe to call once per slide. Leaves real placeholders (title/body/footer/
    slide-number) intact — only deletes the decorative logo + accent strip."""
    layout = slide.slide_layout
    for sh in list(layout.shapes):
        if sh.is_placeholder:
            continue  # never touch placeholders
        L = sh.left/914400 if sh.left is not None else 0
        Tp = sh.top/914400 if sh.top is not None else 0
        W = sh.width/914400 if sh.width else 0
        H = sh.height/914400 if sh.height else 0
        is_logo = (sh.shape_type == _MST.PICTURE and L > 9.5 and Tp < 1.2)
        is_red_strip = (sh.shape_type == _MST.AUTO_SHAPE and W < 0.2 and H > 0.4 and L < 0.6)
        if is_logo or is_red_strip:
            sh._element.getparent().remove(sh._element)

# --- from brand block [4] ---
def center_icons_in_circles(slide, tol_in=0.6, glyph_max_in=1.0):
    """Snap small icon glyphs (groups/freeforms/pictures) to the center of the nearest
    circle/oval within tol_in inches. Idempotent. Returns count moved."""
    ovals = [sh for sh in slide.shapes
             if 'Oval' in (sh.name or '') and sh.width and sh.height]
    if not ovals:
        return 0
    moved = 0
    for g in slide.shapes:
        if g.shape_type not in (_MST2.GROUP, _MST2.FREEFORM, _MST2.PICTURE):
            continue
        if 'Oval' in (g.name or '') or not g.width or not g.height:
            continue
        if g.width / 914400 > glyph_max_in:        # skip large pictures/panels
            continue
        gx = (g.left + g.width / 2) / 914400
        gy = (g.top + g.height / 2) / 914400
        best = None; bd = tol_in
        for o in ovals:
            ox = (o.left + o.width / 2) / 914400
            oy = (o.top + o.height / 2) / 914400
            d = ((gx - ox) ** 2 + (gy - oy) ** 2) ** 0.5
            if d < bd:
                bd = d; best = (ox, oy)
        if best:
            g.left = int(best[0] * 914400 - g.width / 2)
            g.top  = int(best[1] * 914400 - g.height / 2)
            moved += 1
    return moved

# --- from brand block [5] ---
def restore_clipped_widths(out_prs, src_prs, body_offset=1):
    """Undo width-clipping introduced during conversion by restoring each shape's width
    from the matching source shape (matched by name). body_offset aligns output body slides
    to source slides (1 if a title slide was prepended). Returns count restored."""
    src_slides = list(src_prs.slides)
    for oi, out_slide in enumerate(out_prs.slides):
        si = oi - body_offset
        if si < 0 or si >= len(src_slides):
            continue
        src_by_name = {sh.name: sh for sh in src_slides[si].shapes}
        for sh in out_slide.shapes:
            s = src_by_name.get(sh.name)
            if s and s.width and sh.width and sh.width < s.width * 0.9:
                sh.width = s.width

# --- from brand block [6] ---
def center_brand_title(slide, strip_h_in=0.65, band_h_in=0.62):
    """Position the title placeholder centered in the white band just below the strip."""
    for sh in slide.shapes:
        try:
            if sh.is_placeholder and sh.placeholder_format.idx == 0:
                sh.left = Inches(0.663)
                sh.top = Inches(strip_h_in + 0.08)     # clears the 0.65" strip
                sh.width = Inches(11.0)                  # no logo to clear anymore → wider
                sh.height = Inches(band_h_in)
                # vertical-center the text within the placeholder
                bodyPr = sh._element.find('.//' + qn('a:bodyPr'))
                if bodyPr is not None:
                    bodyPr.set('anchor', 'ctr')
                return
        except Exception:
            pass

# --- BRAND CHROME INSERTION (banners + logos, with whitespace-contingent rules) ---
#
# Chrome images live in takeda_templates.py (base64). Each brand's chrome is also
# baked into that brand's *layouts* — so when you build with brand="FRUZAQLA" the
# banner/logo appear automatically via the layout. These helpers are for Mode B
# (reformatting an existing deck onto a clean canvas) where you need to stamp chrome
# onto a slide that does not already carry the brand layout.
#
# Placement geometry is the ground truth extracted from each brand's reference deck
# (see takeda_templates.CHROME_GEOMETRY).

def _chrome_png(name):
    """Return a filesystem path to a decoded brand-chrome image."""
    # Prefer the split-module registry; fall back to a combined takeda_templates.py.
    import builtins as _bi
    reg = getattr(_bi, "_TAKEDA_CHROME_REGISTRY", {})
    if name in reg:
        import base64 as _b64, gzip as _gz
        ext = "jpg" if name == "ENTYVIO_BANNER" else "png"
        cache_dir = "/tmp/takeda_templates"; os.makedirs(cache_dir, exist_ok=True)
        path = os.path.join(cache_dir, f"chrome_{name.lower()}.{ext}")
        if not os.path.exists(path):
            with open(path, "wb") as f:
                f.write(_gz.decompress(_b64.b64decode(reg[name])))
        return path
    import takeda_templates as _tt
    return _tt.chrome_path(name)

def _chrome_geometry():
    """Return the chrome placement geometry from whichever module registered it."""
    import builtins as _bi
    g = getattr(_bi, "_TAKEDA_CHROME_GEOMETRY", None)
    if g:
        return g
    import takeda_templates as _tt
    return _tt.CHROME_GEOMETRY

def _bottom_right_is_clear(slide, zone_left=9.6, zone_top=6.2):
    """Whitespace test for the bottom-right logo zone. Returns True if no content
    shape overlaps the bottom-right region (so a brand logo banner may be placed).
    Footers/slide-number placeholders and full-bleed/full-width chrome are ignored
    (they don't block the logo)."""
    for sh in slide.shapes:
        try:
            if sh.left is None or sh.top is None:
                continue
            L = sh.left / 914400; T = sh.top / 914400
            W = (sh.width or 0) / 914400; H = (sh.height or 0) / 914400
            nm = (sh.name or "").lower()
            if "footer" in nm or "slide number" in nm or "slidenumber" in nm:
                continue
            # ignore full-bleed background / full-width thin chrome strips
            if W > 12.0 and (H > 7.0 or H < 0.5):
                continue
            # does this shape intrude into the bottom-right zone?
            if (L + W) > zone_left and (T + H) > zone_top:
                return False
        except Exception:
            continue
    return True

def _send_behind(slide, pic):
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)
    return pic

def apply_fruzaqla_chrome(slide, place_logo="auto"):
    """FRUZAQLA chrome: top gradient banner (full-bleed image whose top carries the
    gradient + thin border line) sent behind content, PLUS the bottom-right logo
    banner. The logo is whitespace-contingent: with place_logo='auto' it is added
    only when the bottom-right zone is clear of content.
      place_logo: 'auto' (default) | True (force) | False (never)
    """
    geo = _chrome_geometry()
    # Decide on the logo BEFORE adding the full-bleed banner (the banner itself would
    # otherwise register as occupying the bottom-right zone).
    add_logo = place_logo is True or (place_logo == "auto" and _bottom_right_is_clear(slide))
    g = geo["FRUQ_BANNER"]
    banner = slide.shapes.add_picture(
        _chrome_png("FRUQ_BANNER"),
        Inches(g["L"]), Inches(g["T"]), Inches(g["W"]), Inches(g["H"]))
    _send_behind(slide, banner)
    if add_logo:
        gl = geo["FRUQ_LOGO"]
        slide.shapes.add_picture(
            _chrome_png("FRUQ_LOGO"),
            Inches(gl["L"]), Inches(gl["T"]), Inches(gl["W"]), Inches(gl["H"]))
    return banner

def apply_entyvio_chrome(slide):
    """ENTYVIO chrome: bottom full-width gradient banner (purple→crimson→red strip),
    sent behind content. Always applied (it's a footer strip, not whitespace-contingent)."""
    geo = _chrome_geometry()
    g = geo["ENTYVIO_BANNER"]
    strip = slide.shapes.add_picture(
        _chrome_png("ENTYVIO_BANNER"),
        Inches(g["L"]), Inches(g["T"]), Inches(g["W"]), Inches(g["H"]))
    return _send_behind(slide, strip)

def apply_iclusig_chrome(slide, width=None, place_logo=None, **_legacy):
    """Place the ICLUSIG bottom-right logo on ONE slide, anchored in the corner
    just above the footer. `width` (inches) sets the size — pass the value returned
    by iclusig_logo_fit() so every slide in the deck shares one uniform size that
    clears all content. If width is None, the module default (small) is used.

    NOTE: this is the per-slide primitive. For a whole deck use stamp_iclusig_chrome(),
    which also bakes the top line into the layouts, sizes the logo to fit every slide,
    and skips the cover + dividers automatically."""
    if place_logo is False:               # legacy: explicit skip
        return None
    geo = _chrome_geometry()["ICLUSIG_LOGO"]
    w = float(width) if width else geo["W"]
    h = w / geo.get("ASPECT", 1186 / 511)
    L = geo.get("ANCHOR_RIGHT", 13.30) - w
    T = geo.get("ANCHOR_BOTTOM", 7.14) - h
    return slide.shapes.add_picture(
        _chrome_png("ICLUSIG_LOGO"), Inches(L), Inches(T), Inches(w), Inches(h))

# ── ICLUSIG top line (fixed template chrome, baked into layouts) ─────────────
def stamp_iclusig_top_line(prs):
    """Bake the ICLUSIG three-segment top line into EVERY slide layout so it shows
    on all slides as fixed chrome (not selectable/movable on a slide). Idempotent:
    removes any existing 'brandline_*' shapes (in layouts AND slides) first, so a
    per-slide line from an earlier pass is replaced by the layout-level one.
    Call once per deck."""
    from pptx.oxml.ns import qn as _qn
    spec = _chrome_geometry()["ICLUSIG_TOP_LINE"]
    W_EMU = int(prs.slide_width); H_EMU = int(spec["H"] * 914400); T_EMU = int(spec["T"] * 914400)

    def _seg_xml(idx, a, b, hexc):
        L = round(a * W_EMU); Wd = round(b * W_EMU) - L
        return (f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<p:nvSpPr><p:cNvPr id="{9000+idx}" name="brandline_{idx}"/>'
                f'<p:cNvSpPr/>'
                f'<p:nvPr/></p:nvSpPr><p:spPr>'
                f'<a:xfrm><a:off x="{L}" y="{T_EMU}"/><a:ext cx="{Wd}" cy="{H_EMU}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'<a:solidFill><a:srgbClr val="{hexc}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></p:spPr>'
                f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')

    def _strip(spTree):
        for sp in list(spTree.findall(_qn("p:sp"))):
            nv = sp.find(_qn("p:nvSpPr") + "/" + _qn("p:cNvPr"))
            if nv is not None and (nv.get("name") or "").startswith("brandline_"):
                spTree.remove(sp)

    # remove any per-slide copies from a previous run
    for sl in prs.slides:
        _strip(sl.shapes._spTree)
    # add to every layout (drawn above the layout background, under slide content)
    n = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            spTree = layout.shapes._spTree
            _strip(spTree)
            for i, (a, b, hexc) in enumerate(spec["SEGMENTS"]):
                spTree.append(parse_xml(_seg_xml(i, a, b, hexc)))
            n += 1
    return n

# ── ICLUSIG logo auto-sizing (fit every content slide, uniform) ─────────────
def iclusig_logo_fit(prs, content_slides=None, min_w=0.6, max_w=2.1, step=0.05):
    """Return the largest uniform logo WIDTH (inches, anchored bottom-right) that
    clears content on every content slide, using shape bounding boxes (no render
    dependency; conservative). `content_slides` = iterable of 1-based indices to
    consider; default = all slides that aren't the cover or a dark divider (caller
    normally passes the set it will actually stamp)."""
    geo = _chrome_geometry()["ICLUSIG_LOGO"]
    Rx = geo.get("ANCHOR_RIGHT", 13.30); By = geo.get("ANCHOR_BOTTOM", 7.14)
    aspect = geo.get("ASPECT", 1186 / 511)
    EMU = 914400
    slides = list(prs.slides)
    idxs = content_slides if content_slides is not None else range(1, len(slides) + 1)

    def _boxes(slide):
        out = []
        for sh in slide.shapes:
            nm = (sh.name or "")
            if nm.startswith("brandline"):
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            L = sh.left / EMU; T = sh.top / EMU; W = sh.width / EMU; H = sh.height / EMU
            if W <= 0 or H <= 0 or L > 13.4 or T > 7.5 or L + W < 0 or T + H < 0:
                continue
            # ignore full-bleed backgrounds and full-width thin strips (chrome)
            if W > 12.0 and (H > 7.0 or H < 0.5):
                continue
            out.append((L, T, L + W, T + H))
        return out

    best = max_w
    for i in idxs:
        if not (1 <= i <= len(slides)):
            continue
        boxes = _boxes(slides[i - 1])
        w = max_w
        while w >= min_w:
            h = w / aspect
            x0, y0 = Rx - w, By - h
            if not any(a < Rx and c > x0 and b < By and d > y0 for (a, b, c, d) in boxes):
                break
            w -= step
        best = min(best, max(w, min_w))
    return round(best, 2)

def recolor_iclusig_dividers(prs, divider_layout_names=("Section Divider Dark",),
                             hexc="1077BE"):
    """Set the background of ICLUSIG section-divider layouts to solid brand blue
    (#1077BE) — the brand rule is a blue section slide, not purple/teal. Operates
    on the layout background so all divider slides update at once."""
    from pptx.oxml.ns import qn as _qn
    changed = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name not in divider_layout_names:
                continue
            cSld = layout._element.find(_qn("p:cSld"))
            if cSld is None:
                continue
            bg = cSld.find(_qn("p:bg"))
            if bg is not None:
                cSld.remove(bg)
            bg_xml = (f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                      f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                      f'<p:bgPr><a:solidFill><a:srgbClr val="{hexc}"/></a:solidFill>'
                      f'<a:effectLst/></p:bgPr></p:bg>')
            cSld.insert(0, parse_xml(bg_xml))
            changed += 1
    return changed

def recolor_iclusig_subheads(prs, hexc="FCD206"):
    """Recolor sub-heading TEXT (theme accent2 runs) to ICLUSIG yellow so heading
    (blue) and sub-heading are visually distinct. Only touches run text — shape
    fills that use accent2 (e.g. table headers) are left alone."""
    from pptx.enum.dml import MSO_THEME_COLOR as _TC
    n = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None and \
                           run.font.color.theme_color == _TC.ACCENT_2:
                            run.font.color.rgb = RGBColor.from_string(hexc); n += 1
                    except Exception:
                        pass
    return n

def place_iclusig_logo_on_layouts(prs, layout_names, width, logo_path=None):
    """Add the ICLUSIG logo as FIXED chrome on the given layouts (by name) so it
    shows on every slide using them, appears in Slide Master view, and is NOT
    movable on a slide. Removes any pre-existing per-slide logo copies first.
    `width` inches (anchored bottom-right, aspect from geometry). Returns count."""
    geo = _chrome_geometry()["ICLUSIG_LOGO"]
    aspect = geo.get("ASPECT", 1186 / 511)
    Rx = geo.get("ANCHOR_RIGHT", 13.30); By = geo.get("ANCHOR_BOTTOM", 7.14)
    w = float(width); h = w / aspect
    L = round((Rx - w) * 914400); T = round((By - h) * 914400)
    W = round(w * 914400); H = round(h * 914400)
    path = logo_path or _chrome_png("ICLUSIG_LOGO")

    # strip any per-slide logo copies (bottom-right pictures named ICLUSIG_LOGO or
    # any picture parked in the bottom-right logo zone)
    for sl in prs.slides:
        for sh in list(sl.shapes):
            nm = (sh.name or "")
            if nm == "ICLUSIG_LOGO" or (
                sh.shape_type == _MST.PICTURE and sh.left is not None and sh.top is not None
                and sh.left / 914400 > 11 and sh.top / 914400 > 6):
                sh._element.getparent().remove(sh._element)

    names = set(layout_names)
    n = 0
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name not in names:
                continue
            # drop an existing logo on this layout (idempotent)
            for sh in list(lay.shapes):
                if (sh.name or "") == "ICLUSIG_LOGO":
                    sh._element.getparent().remove(sh._element)
            _ip, rid = lay.part.get_or_add_image_part(path)
            pic = (f'<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                   f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                   f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                   f'<p:nvPicPr><p:cNvPr id="{9600+n}" name="ICLUSIG_LOGO"/>'
                   f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/>'
                   f'</p:cNvPicPr><p:nvPr/></p:nvPicPr>'
                   f'<p:blipFill><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
                   f'<p:spPr><a:xfrm><a:off x="{L}" y="{T}"/><a:ext cx="{W}" cy="{H}"/></a:xfrm>'
                   f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>')
            lay.shapes._spTree.append(parse_xml(pic))
            n += 1
    return n

def stamp_iclusig_chrome(prs, cover_slides=(1,), divider_layout_names=("Section Divider Dark",),
                         do_dividers=True, do_subheads=True, logo_layout_names=None):
    """One call to apply the full ICLUSIG chrome/branding to a finished deck — all as
    FIXED template chrome editable in Slide Master view (nothing per-slide/movable):
      1. bake the top line into all layouts (shows on every slide incl. cover/dividers,
         which hide MASTER graphics — hence layouts, not the master);
      2. recolor divider layouts to brand blue (optional);
      3. recolor sub-heading text to brand yellow (optional);
      4. size the logo to fit every content slide (uniform) and place it on the
         CONTENT LAYOUTS — so it shows on content slides and is auto-excluded from the
         cover + dividers (which don't carry it).
    `logo_layout_names`: override the content-layout set; default = the layouts used by
    slides that are neither the cover nor a divider. Call AFTER content + titles are set."""
    stamp_iclusig_top_line(prs)
    if do_dividers:
        recolor_iclusig_dividers(prs, divider_layout_names)
    subs = recolor_iclusig_subheads(prs) if do_subheads else 0

    divider_idx = set()
    for i, sl in enumerate(prs.slides, 1):
        try:
            if sl.slide_layout.name in divider_layout_names:
                divider_idx.add(i)
        except Exception:
            pass
    skip = set(cover_slides) | divider_idx
    content = [i for i in range(1, len(prs.slides) + 1) if i not in skip]
    w = iclusig_logo_fit(prs, content_slides=content)

    # content layouts = layouts used by the content slides (unless overridden)
    if logo_layout_names is None:
        slides = list(prs.slides)
        logo_layout_names = sorted({slides[i - 1].slide_layout.name for i in content})
    nlogo = place_iclusig_logo_on_layouts(prs, logo_layout_names, w)
    return {"logo_width_in": w, "logo_layouts": nlogo, "logo_layout_names": list(logo_layout_names),
            "subheads_recolored": subs, "dividers_blue": do_dividers}

def apply_brand_chrome(slide, brand, **kw):
    """Dispatch to the right brand chrome function."""
    b = (brand or "").upper()
    if b == "FRUZAQLA": return apply_fruzaqla_chrome(slide, **kw)
    if b == "ENTYVIO":  return apply_entyvio_chrome(slide)
    if b == "ICLUSIG":  return apply_iclusig_chrome(slide, **kw)
    return None

# Back-compat alias (old name).
def apply_fruzaqla_banner(slide):
    """Deprecated: use apply_fruzaqla_chrome(). Kept for back-compat."""
    return apply_fruzaqla_chrome(slide, place_logo="auto")

# --- from brand block [16] ---
def audit_mode_b_fidelity(src_prs, out_prs):
    """Flag slides where conversion lost content or rasterized the source.
    Returns a list of (slide_index, problem) — empty list means PASS."""
    problems = []
    src_slides = list(src_prs.slides)
    out_slides = [s for s in out_prs.slides]  # skip title/end if offset; align by content
    for i, (ss, os_) in enumerate(zip(src_slides, out_slides)):
        # 1. Full-slide raster = the rasterization signature → always a failure.
        #    Distinguish from legitimate large content images: a raster artifact is
        #    full-BLEED at the origin (≈0,0) covering essentially the whole 13.33×7.5"
        #    canvas. A real content image is inset (top > ~0.5" or not full width).
        for sh in os_.shapes:
            if sh.shape_type == _T.PICTURE and sh.width and sh.height and sh.left is not None and sh.top is not None:
                L, Tp = sh.left/914400, sh.top/914400
                W, H = sh.width/914400, sh.height/914400
                if L <= 0.2 and Tp <= 0.2 and W >= 13.0 and H >= 7.2:
                    problems.append((i+1, "FULL-SLIDE RASTER — source was flattened to an image"))
                    break
        # 2. Charts/tables must survive the conversion
        src_charts = sum(1 for sh in ss.shapes if sh.has_chart)
        out_charts = sum(1 for sh in os_.shapes if sh.has_chart)
        if out_charts < src_charts:
            problems.append((i+1, f"LOST CHARTS: {src_charts} → {out_charts}"))
        src_tables = sum(1 for sh in ss.shapes if sh.has_table)
        out_tables = sum(1 for sh in os_.shapes if sh.has_table)
        if out_tables < src_tables:
            problems.append((i+1, f"LOST TABLES: {src_tables} → {out_tables}"))
        # 3. Gross shape-count collapse (allow +/- a few for chrome adjustments)
        if len(ss.shapes) >= 5 and len(os_.shapes) < len(ss.shapes) * 0.5:
            problems.append((i+1, f"SHAPE COLLAPSE: {len(ss.shapes)} → {len(os_.shapes)} (content likely lost)"))
        # 4. Leftover Takeda corporate chrome (logo / red accent strip) on a brand slide.
        #    These live on the layout; neutralize_takeda_layout_chrome() should have removed them.
        for sh in os_.shapes:
            try:
                L = sh.left/914400 if sh.left is not None else 0
                W = sh.width/914400 if sh.width else 0
                if sh.shape_type == _T.PICTURE and L > 9.5 and (sh.top or 0)/914400 < 1.2:
                    # could be the FRUQ strip (full width) — only flag if NOT full width
                    if W < 6.0:
                        problems.append((i+1, "TAKEDA LOGO still present (layout chrome not neutralized)"))
            except Exception:
                pass
    # 5. Brand-accent presence (deck-level, not per-slide): if the output uses essentially
    #    no brand accent colors, the remap dict almost certainly didn't cover the source
    #    palette and everything collapsed to gray/black via the catch-all. This is the
    #    "flat, accent-less deck" failure the ZS-format conversion hit.
    #    IMPORTANT: count BOTH literal brand srgbClr AND theme-resolved schemeClr refs.
    #    After retheme_brand(), accents live in the theme and are referenced via schemeClr,
    #    so a literal-only count reads ~0 and false-positives. Count both.
    accent_set = {'9A4398','743272','074F71','BEB7D7','EDD7ED','81C44D'}  # FRUQ accents
    accent_scheme = {'accent1','accent2','accent3','accent4','accent5','accent6'}
    accent_hits = 0
    for slide in out_prs.slides:
        for e in slide._element.iter(qn('a:srgbClr')):
            if e.get('val','').upper() in accent_set:
                accent_hits += 1
        for e in slide._element.iter(qn('a:schemeClr')):
            if e.get('val') in accent_scheme:
                accent_hits += 1
    if accent_hits < max(3, len(list(out_prs.slides))):  # ~<1 accent ref per slide → suspicious
        problems.append((0, f"ACCENT-LESS OUTPUT: only {accent_hits} brand-accent refs "
                            f"(srgb + scheme) across the deck — either BRAND_REMAP does not "
                            f"cover the source palette, or retheme_brand() was not run."))
    # 6. Residual Takeda corporate red anywhere = theme/chevron/chart not rethemed.
    #    After retheme_brand(), accent1 (E1242A) must be gone; if it survives as a literal
    #    srgbClr OR the theme still defines it, schemeClr-driven shapes will render red.
    red_refs = sum(1 for slide in out_prs.slides
                   for e in slide._element.iter(qn('a:srgbClr'))
                   if e.get('val','').upper() == 'E1242A')
    if red_refs:
        problems.append((0, f"TAKEDA RED LEAK: {red_refs} literal E1242A refs remain — extend "
                            f"BRAND_REMAP. (Also confirm retheme_brand() ran for schemeClr shapes.)"))
    return problems

def _count_color(prs, hex6):
    """Count literal srgbClr refs to a hex value across all slides (QA helper)."""
    hex6 = hex6.upper()
    return sum(1 for slide in prs.slides
               for e in slide._element.iter(qn('a:srgbClr'))
               if e.get('val','').upper() == hex6)

# --- from brand block [29] ---
def _e(v): return str(int(round(v * EMU)))

def apply_entyvio_chrome_to_master(prs):
    """Place the Takeda blade + ENTYVIO footer strip on the slide master as LOCKED,
    non-selectable furniture. Call ONCE after the deck is built. Removes any per-slide
    footer pictures first so chrome isn't duplicated."""
    master = prs.slide_masters[0]
    mtree = master.element.find(qn('p:cSld')).find(qn('p:spTree'))

    # 1. footer strip image — use the EDGE-BLEED variant (gray margins cropped off the
    #    native PNG) so the gradient reaches the slide perimeter with no white gaps.
    #    See "Footer Strip Asset" note below: crop the ~#F0F0F0 margins before embedding.
    with open(_get_entyvio_strip_bleed(), 'rb') as f: strip_blob = f.read()
    img_part, rId = master.part.get_or_add_image_part(io.BytesIO(strip_blob))

    # 2. locked footer picture (bottom 0.35", behind content) — FULL WIDTH 0 -> 13.333"
    #    The image MUST be the bled version; a rect prstGeom at full slide width then
    #    shows gradient edge-to-edge.
    foot = f'''<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">
      <p:nvPicPr>
        <p:cNvPr id="9001" name="entyvio_footer_strip"/>
        <p:cNvPicPr><a:picLocks noSelect="1" noMove="1" noResize="1"/></p:cNvPicPr>
        <p:nvPr userDrawn="1"/>
      </p:nvPicPr>
      <p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
      <p:spPr><a:xfrm><a:off x="0" y="{_e(7.15)}"/><a:ext cx="{_e(13.333)}" cy="{_e(0.35)}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
    </p:pic>'''
    mtree.append(etree.fromstring(foot))

    # 3. locked AUTHENTIC Takeda blade — exact custGeom path lifted from
    #    Takeda_Slide_Template_EN.potx ("Rectangle 1" on the title layouts).
    #    Path is a curved petal: straight left edge, flat top/bottom stubs, cubic Bézier
    #    bulge on the right. Path coordinate space is w=190787 h=565121 (do NOT rescale the
    #    path numbers — only the <a:ext> sets the on-slide size). The rect uses NAMED guides
    #    (l/t/r/b) exactly as the template does; numeric bounds would clip the Bézier control
    #    points (235226 > path width) and is the bug that previously broke the file.
    #    Fill is Takeda red #E1242A (NOT accent1, which is ENTYVIO purple) — the blade is the
    #    one deliberate Takeda-red element on the deck.
    pw, ph = int(BLADE['W']*EMU), int(BLADE['H']*EMU)
    blade = f'''<p:sp xmlns:p="{P}" xmlns:a="{A}">
      <p:nvSpPr>
        <p:cNvPr id="9002" name="takeda_header_blade"/>
        <p:cNvSpPr><a:spLocks noSelect="1" noMove="1" noResize="1"/></p:cNvSpPr>
        <p:nvPr userDrawn="1"/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{_e(BLADE['L'])}" y="{_e(BLADE['T'])}"/><a:ext cx="{pw}" cy="{ph}"/></a:xfrm>
        <a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
          <a:rect l="l" t="t" r="r" b="b"/>
          <a:pathLst><a:path w="190787" h="565121">
            <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
            <a:lnTo><a:pt x="66951" y="0"/></a:lnTo>
            <a:cubicBezTo>
              <a:pt x="235226" y="153449"/>
              <a:pt x="228876" y="408497"/>
              <a:pt x="66951" y="565121"/>
            </a:cubicBezTo>
            <a:lnTo><a:pt x="0" y="565121"/></a:lnTo>
            <a:lnTo><a:pt x="0" y="0"/></a:lnTo>
            <a:close/>
          </a:path></a:pathLst>
        </a:custGeom>
        <a:solidFill><a:srgbClr val="{BLADE['fill']}"/></a:solidFill>
        <a:ln w="28575"><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr wrap="square" rtlCol="0" anchor="ctr"/><a:lstStyle/><a:p><a:pPr algn="ctr"/><a:endParaRPr lang="en-US"/></a:p></p:txBody>
    </p:sp>'''
    mtree.append(etree.fromstring(blade))

    # 4. remove any per-slide footer pictures (avoid duplicate clickable chrome)
    for s in prs.slides:
        for sh in list(s.shapes):
            if sh.shape_type == 13 and sh.top is not None and abs(sh.top/EMU - 7.15) < 0.06:
                sh._element.getparent().remove(sh._element)

# --- from brand block [30] ---
def apply_entyvio_footer(slide, slide_height_in=7.5, strip_h_in=0.35):
    """DEPRECATED (v0.7) — per-slide footer picture. Use apply_entyvio_chrome_to_master(prs)."""
    from pptx.util import Inches
    pic = slide.shapes.add_picture(
        _get_entyvio_strip(),
        Inches(0.0), Inches(slide_height_in - strip_h_in),
        Inches(13.333), Inches(strip_h_in),
    )
    spTree = slide.shapes._spTree
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)
    return pic

# --- from brand block [31] ---
def _get_entyvio_strip():
    """Decode the embedded ENTYVIO footer strip to a temp PNG; return its path."""
    cache = "/tmp/entyvio_footer_strip.png"
    if not _os.path.exists(cache):
        with open(cache, "wb") as f:
            f.write(_b64.b64decode(ENTYVIO_STRIP_B64))
    return cache

# --- from brand block [3] (retheme_brand; FRUQ_THEME_REMAP moved to takeda_brands) ---
from lxml import etree
_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
def retheme_brand(prs, theme_remap):
    """Rewrite the presentation theme's accent palette to brand colors so EVERY schemeClr
    reference (charts, chevrons, SmartArt, nested icon groups) resolves to the brand palette.
    Deterministic, one pass per deck. Returns count of accent slots rewritten.
    NOTE: only remaps accent1-6. Leaves dk1/lt1/dk2/lt2 (text/background slots) alone so
    body text stays readable. Theme parts are generic Parts (use .blob, not ._element)."""
    n = 0; seen = set()
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype and rel.target_part.partname not in seen:
                seen.add(rel.target_part.partname)
                tp = rel.target_part
                root = etree.fromstring(tp.blob)
                cs = root.find(f'.//{{{_A}}}clrScheme')
                if cs is None:
                    continue
                for acc, hexv in theme_remap.items():
                    node = cs.find(f'{{{_A}}}{acc}')
                    if node is None:
                        continue
                    srgb = node.find(f'{{{_A}}}srgbClr')
                    sysc = node.find(f'{{{_A}}}sysClr')
                    if srgb is not None:
                        srgb.set('val', hexv); n += 1
                    elif sysc is not None:           # convert sysClr slot → srgbClr
                        new = etree.SubElement(node, f'{{{_A}}}srgbClr'); new.set('val', hexv)
                        node.remove(sysc); n += 1
                tp._blob = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
    return n

# --- from brand block [7] (function only; example wrapper dropped) ---
def strip_header_bar_fills(slide, top_limit_in=1.20):
    """Remove solid fills from wide shapes in the title zone so no colored box remains."""
    for sh in list(slide.shapes):
        try:
            if sh.top is None: continue
            t = sh.top / 914400
            w = sh.width / 914400 if sh.width else 0
            h = sh.height / 914400 if sh.height else 0
            # Wide, short shape sitting in the header zone = likely a header bar
            if t < top_limit_in and w > 6 and h < 1.2:
                spPr = sh._element.find('.//' + qn('p:spPr'))
                if spPr is not None and spPr.find(qn('a:solidFill')) is not None:
                    # If it has no meaningful text, delete it; else just remove the fill
                    has_text = sh.has_text_frame and sh.text_frame.text.strip()
                    if not has_text:
                        sh._element.getparent().remove(sh._element)
                    else:
                        for tag in ('a:solidFill','a:gradFill','a:blipFill','a:pattFill'):
                            f = spPr.find(qn(tag))
                            if f is not None: spPr.remove(f)
                        if spPr.find(qn('a:noFill')) is None:
                            spPr.append(spPr.makeelement(qn('a:noFill'), {}))
        except Exception:
            pass


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — PER-BRAND COLOR REGISTRY
# Usage: b = get_brand("FRUZAQLA"); remap = b["BRAND_REMAP"]
# ═══════════════════════════════════════════════════════════════════════════════
_SRC = {}

_SRC['FRUZAQLA'] = r"""
BRAND_PRIMARY   = '074F71'    # Dark Teal — header TEXT color (theme accent3); pairs with banner
BRAND_ACCENT    = '9A4398'    # Purple — primary accent, callouts, highlights
BRAND_DARK_ACC  = '743272'    # Dark Purple — table sub-headers, deep emphasis
BRAND_LIGHT_ACC = 'BEB7D7'    # Light Purple — table category fills, light emphasis
BRAND_PALE_ACC  = 'EDD7ED'    # Pale Purple — subtle fills, alternating row tint
BRAND_LIGHT     = 'F4F3F3'    # Near-White — card fills, light backgrounds
BRAND_MID_GRAY  = 'A3A1A8'    # Mid Gray — supporting text, dividers
BRAND_DARK      = '1A1B1F'    # Near-Black — primary body text
BRAND_FONT      = 'Aptos'     # Brand font (standardized to Aptos; source deck used Arial)
BRAND_STRIP     = 'fruzaqla_header_strip.png'  # Header STRIP (top 0.65" only — never full slide)

BRAND_PALETTE_SET = {
    '074F71','9A4398','743272','BEB7D7','EDD7ED',
    '1A1B1F','34373F','A3A1A8','F4F3F3','FFFFFF','000000',
    '00B050','81C44D','FF0000','C00000',
    '1A1628','074F71','767171','75737D'
}

BRAND_REMAP = {
    # ZS source oranges → Purple accent
    'EC7200': '9A4398', 'F7941D': '9A4398', 'FF6600': '9A4398',
    'FFA500': '9A4398', 'FFD700': 'BEB7D7',
    # Source blues → Dark Teal
    '0070C0': '074F71', '4472C4': '074F71', '5B9BD5': 'BEB7D7',
    '1077BE': '074F71', '1177BD': '074F71', '01AFD7': '9A4398',
    # Source darks → Near-Black
    '333333': '1A1B1F', '404040': '1A1B1F',
    # Source teals → keep or map to brand teal
    '32A29B': '074F71', '009999': '074F71',
    # Source light grays → Light Gray
    'D0D0D0': 'F4F3F3', 'D8D8D8': 'F4F3F3', 'E8E8E8': 'F4F3F3',
    'F2F2F2': 'F4F3F3', 'EEEEEE': 'F4F3F3', 'E7E6E6': 'F4F3F3',
    # Takeda corporate red → Purple accent (if converting from Takeda corporate)
    'E1242A': '9A4398', '891515': '743272',
    # Source mid grays
    'A1A4AC': 'A3A1A8', 'BFBFBF': 'A3A1A8',

    # ── ZS-FORMAT SOURCE PALETTE (the "State of KAM in Life Sciences" deck family) ──
    # CRITICAL: ZS-authored decks do NOT use Office defaults. Their workhorse text/shape
    # color is 53565A (a warm charcoal-gray, ~239 uses in a 10-slide deck), with blue
    # 00629B, orange ED8B00, teal 4F868E, and plum 6E2B62 accents. None of these were
    # covered before this version, so they ALL fell through to the luminance catch-all
    # and flattened to near-black/gray — wiping out every brand accent. Map them explicitly.
    '53565A': '34373F', '63666A': '34373F', '888B8D': 'A3A1A8',  # ZS grays → charcoal/gray
    '96999E': 'A3A1A8', '6D6E71': '75737D',
    '00629B': '074F71', '004B87': '074F71', '0033A0': '074F71',  # ZS blues → dark teal
    'DDF3FF': 'BEB7D7', 'A5D8F3': 'BEB7D7', 'C6E6F5': 'EDD7ED',  # ZS pale blues → light purple
    'ED8B00': '9A4398', 'FFA300': '9A4398', 'FBB034': 'BEB7D7',  # ZS oranges → purple accent
    '4F868E': '074F71', '72A3AA': 'BEB7D7', '008C95': '074F71',  # ZS teals → dark teal
    '6E2B62': '743272', '8C2C7E': '743272',                       # ZS plums → dark purple
    'B3CA89': '81C44D', '6CC24A': '81C44D',                       # ZS greens → brand green
}

FRUZAQLA_PALETTE = {
    '074F71','9A4398','743272','BEB7D7','EDD7ED',
    '1A1B1F','1A1628','34373F','A3A1A8','75737D','767171',
    'F4F3F3','FFFFFF','000000',
    '00B050','81C44D','FF0000','C00000','008000',
    # Data viz extended (from reference deck)
    'D9E8FF','8CBAFF','B6E4FA','81C7D1',
    'FFFFA7','FFDF85','7FD07D','D4E8C6',
    '7030A0','002060','C06EBE','DCAEDB','E7C7E6',
    'D0CECE','7F7F7F','A6A6A6','BFBFBF','F2F2F2',
    '44546A','0D0D0D','F1F1F1',
}

FRUQ_THEME_REMAP = {
    'accent1': '9A4398',   # was E1242A (Takeda red)  → FRUZAQLA purple
    'accent2': '743272',   # was 891515 (dark red)    → dark purple
    'accent3': '074F71',   # was 34373F (charcoal)    → dark teal
    'accent4': 'BEB7D7',   # was A1A4AC (gray)         → light purple
    'accent5': 'EDD7ED',   # was EDF2F4 (pale)         → pale purple
    'accent6': 'A3A1A8',   #                           → brand gray
}
"""
_SRC['ICLUSIG'] = r"""
BRAND_PRIMARY   = '1077BE'    # Blue — header TITLE TEXT color (no-fill text box, never a bar)
BRAND_ACCENT    = '01AFD7'    # Cyan/Teal — table headers, secondary accent
BRAND_DARK_ACC  = '074F71'    # Dark Teal — deep emphasis, sub-headers
BRAND_DARK_GRN  = '155D39'    # Dark Green — tertiary accent
BRAND_LIGHT     = 'F4F3F3'    # Near-White — card fills, light backgrounds
BRAND_MID_GRAY  = '6F7073'    # Mid Gray — supporting text, dividers
BRAND_DARK      = '1A1B1F'    # Near-Black — primary body text
BRAND_FONT      = 'Aptos'     # Brand font (standardized to Aptos; source deck used Arial)
# Note: #1177BD from source decks is standardized to #1077BE (single brand blue)

BRAND_PALETTE_SET = {
    '1077BE','01AFD7','074F71','155D39',
    'EE1100','1A1B1F','34373F','6F7073','454545',
    'F4F3F3','FFFFFF','000000',
    '00B050','81C44D','FF0000','C00000',
    '1A1628','9A4398','743272','FFE697',
    'DDDDDD','E7E6E6','F2F2F2','44546A',
}

BRAND_REMAP = {
    # ZS source oranges → Cyan accent
    'EC7200': '01AFD7', 'F7941D': '01AFD7', 'FF6600': '01AFD7',
    'FFA500': '01AFD7', 'FFD700': 'FFE697',
    # Source generic blues → Brand blue
    '0070C0': '1077BE', '4472C4': '1077BE', '5B9BD5': '01AFD7',
    # Standardize the two ICLUSIG blues → single brand blue
    '1177BD': '1077BE',
    # Source darks → Near-Black
    '333333': '1A1B1F', '404040': '1A1B1F',
    # Source teals → Brand teal
    '32A29B': '074F71', '009999': '074F71',
    # Source light grays → Light Gray
    'D0D0D0': 'F4F3F3', 'D8D8D8': 'F4F3F3', 'E8E8E8': 'F4F3F3',
    'EEEEEE': 'F4F3F3',
    # Takeda corporate red → Brand blue (if converting from Takeda corporate)
    'E1242A': '1077BE', '891515': '074F71',
    # Source mid grays
    'A1A4AC': '6F7073', 'BFBFBF': '6F7073',
    # Source purples keep (shared ZS)
    # Source greens keep (shared ZS)
}

ICLUSIG_PALETTE = {
    '1077BE','01AFD7','074F71','155D39',
    'EE1100','1A1B1F','1A1628','34373F','6F7073','454545',
    'F4F3F3','F2F2F2','FFFFFF','000000',
    '00B050','81C44D','FF0000','C00000',
    '9A4398','743272','FFE697','DDDDDD','E7E6E6',
    '44546A','0D0D0D',
    # Extended data viz (from reference deck)
    'B9DFF9','80C5F4','138CDF','3AA5EE',
    '3CA2A6','51B1BF','B6E4FA',
    'FCD206','FBD206','FEF2B4',
    'FF6E63','E97132',
    'CBD5D2','D9D9D9','7F7F7F','A6A6A6','BFBFBF',
    '5B9BD5','F1F1F1',
}

"""
_SRC['ALUNBRIG'] = r"""
BRAND_PRIMARY   = '00A1B0'    # Teal — header TITLE TEXT (no-fill text box), chart primary
BRAND_DARK_ACC  = '005158'    # Dark Teal — deep emphasis, sub-title text
BRAND_LIGHT     = 'D4D6DB'    # Silver Gray — table alternating rows, card fills
BRAND_MID_GRAY  = 'A1A4AC'    # Mid Gray — supporting text, dividers
BRAND_DARK      = '34373F'    # Charcoal — secondary body text, dark table rows
BRAND_FONT      = 'Aptos'     # Brand font (standardized to Aptos; source deck used Corbel)

BRAND_PALETTE_SET = {
    '00A1B0','00A0AF','0BA5B3','005158',
    '34373F','4C4948','A1A4AC','D4D6DB','D9D9D9',
    'FFFFFF','000000',
    '00B050','FF0000','C00000',
    '0070C0','7030A0','B9BF03','FFC000',
    'FFFF00','F79646','1F497D','8B8B8B',
    '7F7F7F','F2F2F2','262626',
}

BRAND_REMAP = {
    # Source oranges → Gold
    'EC7200': 'FFC000', 'F7941D': 'FFC000', 'FF6600': 'FFC000',
    'FFA500': 'FFC000',
    # Source blues → Brand teal
    '4472C4': '00A1B0', '5B9BD5': '00A1B0',
    '1077BE': '00A1B0', '1177BD': '00A1B0', '01AFD7': '00A1B0',
    # Source teals → Brand teal (normalize variants)
    '00A0AF': '00A1B0', '0BA5B3': '00A1B0',
    '32A29B': '00A1B0', '009999': '00A1B0',
    # Source purples → keep (chart series)
    '9A4398': '7030A0', '743272': '7030A0',
    # Source darks → Charcoal
    '333333': '34373F', '404040': '34373F', '1A1B1F': '000000',
    # Source light grays → Silver Gray
    'D0D0D0': 'D4D6DB', 'D8D8D8': 'D4D6DB', 'E8E8E8': 'D4D6DB',
    'F2F2F2': 'D9D9D9', 'EEEEEE': 'D9D9D9', 'E7E6E6': 'D9D9D9',
    'F4F3F3': 'D4D6DB',
    # Takeda corporate red → Brand teal (if converting from Takeda corporate)
    'E1242A': '00A1B0', '891515': '005158',
    # Source mid grays → Mid Gray
    'BFBFBF': 'A1A4AC', 'A6A6A6': 'A1A4AC',
    # ICLUSIG/FRUZAQLA brand colors → ALUNBRIG equivalents
    '074F71': '005158',
}

ALUNBRIG_PALETTE = {
    '00A1B0','00A0AF','0BA5B3','005158',
    '34373F','4C4948','A1A4AC','D4D6DB','D9D9D9',
    'FFFFFF','000000',
    '00B050','FF0000','C00000','FFFF00',
    '0070C0','7030A0','B9BF03','FFC000','FFC001',
    'F79646','DC5F12','1F497D','00B0F0',
    '7C41A8','8B8B8B','7F7F7F','F2F2F2','262626',
}

"""
_SRC['ENTYVIO'] = r"""
BRAND_PRIMARY   = '9F1897'    # ENTYVIO Purple — primary brand accent (section headers, key fills)
BRAND_ACCENT    = '9C268F'    # Purple (kicker variant) — top-right kicker tags, sub-headers
BRAND_DARK_ACC  = '9C1E96'    # Deep Purple — deep emphasis, dark purple fills
BRAND_LIGHT_ACC = 'F4D0F0'    # Light Purple — table category fills, light emphasis
BRAND_PALE_ACC  = 'F9DBF8'    # Pale Purple — subtle fills, alternating row tint
BRAND_CRIMSON   = 'B1254A'    # Crimson — title accent tick, KPI percentages, secondary accent
BRAND_RED       = 'E63C2D'    # Bright Red — divider lines, negative semantic
BRAND_ORANGE    = 'EC7200'    # Takeda Orange — tertiary accent (use sparingly)
BRAND_AMBER     = 'FFC000'    # Amber — chart highlight / caution
BRAND_GREEN     = '00B050'    # Green — positive semantic
BRAND_CARD_DARK = '34373F'    # Charcoal — left icon-label cards, dark panels
BRAND_DARKEST   = '1A1628'    # Near-Black (warm) — primary body text
BRAND_MID_GRAY  = 'A3A1A8'    # Mid Gray — supporting text, dividers
BRAND_LIGHT     = 'EDF2F4'    # Pale Blue-Gray — card fills, light backgrounds
BRAND_FONT      = 'Aptos'     # Brand font (standardized to Aptos; Fact Book used Calibri)
BRAND_STRIP     = 'entyvio_footer_strip.png'  # Footer STRIP (bottom 0.35" only — never full slide)
BRAND_STRIP_POS = 'bottom'    # ENTYVIO chrome is a BOTTOM strip (unlike FRUZAQLA's top strip)

BRAND_PALETTE_SET = {
    '9F1897','9C268F','9C1E96','9E1C96','9B288C','DC1A89','D921D0','CC32AF',  # purples
    'F4D0F0','F9DBF8','F5D7EF','F6D6DE',                                       # light/pale purples & pinks
    'B1254A','B1156E','E63C2D',                                               # crimson / red
    'EC7200','FFC000','FFDF69','FECF5A',                                      # orange / amber
    '00B050','54A250',                                                        # greens
    '34373F','1A1628','27292F','4C4948','A3A1A8','75737D','767171',           # darks / grays
    'EDF2F4','ECEDEE','D9E3E7','FFFFFF','FEFEFE','000000',                    # neutrals
}

ENTYVIO_THEME_REMAP = {
    'accent1': '9F1897',   # → ENTYVIO purple (primary)
    'accent2': '9C1E96',   # → deep purple
    'accent3': '34373F',   # → charcoal (text/dark — matches existing accent3)
    'accent4': 'B1254A',   # → crimson (secondary accent)
    'accent5': 'F4D0F0',   # → light purple
    'accent6': 'A3A1A8',   # → mid gray
}

BRAND_REMAP = {
    # Takeda corporate / source reds → ENTYVIO purple family
    'E1242A':'9F1897','891515':'9C1E96','C00000':'9C1E96','FF0000':'E63C2D',
    # Source blues → purple primary
    '0070C0':'9F1897','00629B':'9F1897','004B87':'9C1E96','4472C4':'9F1897',
    '1077BE':'9F1897','01AFD7':'DC1A89',
    # ZS-format source palette (carry over from v0.7 coverage) → ENTYVIO
    '53565A':'34373F','63666A':'34373F','96999E':'A3A1A8','6D6E71':'75737D',
    'ED8B00':'EC7200','F7941D':'EC7200','FFA300':'FFC000',
    '4F868E':'9F1897','72A3AA':'F4D0F0','6E2B62':'9C1E96',
    'DDF3FF':'F9DBF8','A5D8F3':'F4D0F0',
    # Source darks / grays
    '333333':'1A1628','404040':'34373F','A1A4AC':'A3A1A8','BFBFBF':'A3A1A8',
    'D0D0D0':'EDF2F4','D8D8D8':'EDF2F4','E8E8E8':'EDF2F4','F2F2F2':'EDF2F4',
    'E7E6E6':'EDF2F4','EEEEEE':'ECEDEE',
    # Greens preserved (semantic)
    'B3CA89':'00B050','6CC24A':'00B050','92D050':'00B050',
}

"""

def _load(src):
    ns = {}
    exec(src, ns)
    return {k: v for k, v in ns.items() if k[0].isalpha() and k.upper() == k}

BRANDS = {name: _load(src) for name, src in _SRC.items()}

# ── ZS-source augmentation (added 2026-07-29, ZS COGNITIVE → ENTYVIO conversion) ──
# ZS proposal decks use slate-teal 688A92 as a recurring chromatic accent (phase
# ramps, credentials chrome). It is not in any harvested remap, so Mode B reskins
# hard-fail the QA palette check. Map it to each brand's secondary accent so the
# ramp it belongs to lands inside one brand family (recolor rule 7).
_ZS_SLATE_TEAL = {
    'ENTYVIO':  'B1254A',   # crimson — series-3 family member
    'FRUZAQLA': '743272',   # deep plum accent
    'ICLUSIG':  '01AFD7',   # cyan accent
    'ALUNBRIG': '005158',   # deep teal accent
}
for _b, _hex in _ZS_SLATE_TEAL.items():
    if _b in BRANDS:
        BRANDS[_b].setdefault('BRAND_REMAP', {}).setdefault('688A92', _hex)
del _b, _hex


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3b — BRAND EXTENSIONS (full palettes + layout inventories + chrome map)
# Harvested from each brand's gold-standard reference deck (≥3-use color threshold).
# These augment the BRAND_REMAP tables above; they do not replace them.
# ═══════════════════════════════════════════════════════════════════════════════

# Complete color palette observed in each brand's reference deck (fills, text, lines,
# chart series, semantic colors, neutrals). Use as the authoritative "is this color
# on-brand?" set for QA. Primary brand colors are the first few in each BRAND_REMAP.
BRAND_FULL_PALETTE = {
    'FRUZAQLA': {
        '000000','0000FF','002060','0070C0','008000','00B050','01AFD7','06BAB3',
        '074F71','0D0D0D','19B861','1A1628','1A1B1F','1A777A','34373F','368995',
        '3CA2A6','44546A','474554','4C4948','51B1BF','595959','5B9BD5','663300',
        '7030A0','70AD47','743272','75737D','767171','777777','7F7F7F','7FD07D',
        '81C44D','81C7D1','8CBAFF','9A4398','A3A1A8','A6A6A6','A6D7DE','B6E4FA',
        'BEB7D7','BFBFBF','C00000','CC9900','D0CECE','D1D0D4','D3F5E4','D4E8C6',
        'D9D9D9','D9E8FF','E2E2E2','E5F6E5','E7C7E6','E7E6E6','E7F2F8','EC7200',
        'EDD7ED','EE1100','F1F1F1','F26B43','F2F2F2','F4F3F3','FBAE40','FF0000',
        'FF6E63','FFD8D5','FFDF85','FFFF00','FFFFA7','FFFFFF',
    },
    'ICLUSIG': {
        '000000','0000FF','0070C0','00B050','01AFD7','074F71','0D0D0D','1077BE',
        '1177BD','155D39','1A1628','1A1B1F','32BF40','34373F','3AA5EE','3CA2A6',
        '44546A','454545','474554','4C4948','50C85C','51B1BF','595959','5B9BD5',
        '663300','6F7073','70AD47','743272','75737D','7F7F7F','81C44D','83CCEB',
        '891515','9A4398','A1A4AC','A3A1A8','A6A6A6','B2E3B1','B6E4FA','B9DFF9',
        'C00000','CBD5D2','D1D0D4','D3F5E4','D9D9D9','DAF2D0','DDDDDD','E1242A',
        'E7E6E6','E7F2F8','E97132','EC7200','EDF2F4','EE1100','F26B43','F2F2F2',
        'F4F3F3','F6BB00','FBAE40','FBD206','FCD206','FEF2B4','FF0000','FF6E63',
        'FFE697','FFFF00','FFFFFF',
    },
    'ENTYVIO': {
        '000000','001423','001965','002060','004D86','0069AA','0070C0','007635',
        '00B050','00B0F0','0E2841','0F9ED5','156082','196B24','19D3C5','1A1628',
        '27292F','2A918B','336E68','34373F','3F4344','467886','474554','4C4948',
        '4EA72E','54A250','5B8595','69FFAD','6C6F6F','7030A0','750B4D','75737D',
        '7B1376','7D8394','7F7F7F','807CB0','82786F','829696','891515','8A1317',
        '92D050','95908F','96607D','990A91','9B288C','9C1E96','9C268F','9E1C96',
        '9F1897','9FE6FF','A02B93','A1A4AC','A3A1A8','A3A3A3','AC181C','B1156E',
        'B1254A','B2DE82','B381D9','B98CDC','BB9301','BCF9FF','BCFFDB','C7A2E3',
        'CB82C7','CC32AF','D0D0D0','D75DC0','D921D0','D9D9D9','D9E3E7','DA68C4',
        'DB1285','DC1A89','E0E9EC','E1242A','E33B3B','E34848','E3D1F1','E450DD',
        'E63C2D','E64A0E','E65CDF','E8E8E8','E97132','EC7200','EC8585','ECEDEE',
        'EDF2F4','EE1100','F0F0F0','F1912B','F26B43','F2F2F2','F4D0F0','F4F3F3',
        'F5D7EF','F6D6DE','F6DAF3','F9D3D4','F9DBF8','FBAE40','FCEAFB','FCFEFF',
        'FECF5A','FEFEFE','FF0000','FFC000','FFC80A','FFCBCB','FFCCCC','FFD129',
        'FFDF69','FFE89F','FFECBD','FFEDAB','FFFF00','FFFFFF',
    },
}

# Deduplicated layout inventory per brand (Option B-trimmed). These are the layout
# NAMES present in each brand template (takeda_templates.template_path(brand)).
# Resolve a layout by name via Deck(brand=...).add_slide_by_name() — never by index.
BRAND_LAYOUTS = {
    'FRUZAQLA': [
        'Title Slide', 'Title_Eyebrow_subhead_Content', '1_Title_Eyebrow_subhead', 'Section Header',
        'Two Content', 'Comparison', 'Blank', 'Advanced Chart 3 Column',
        '3_Title', 'Standard Title Only', 'D. Title Only', 'Title w/ Subhead & Headline',
        'Title Only', 'One Third LHS Dark', 'Standard 2-Column Text', '3_Title and Content',
        'Title Page 1', '1_Title Page 2 ', 'Section Title Page', 'Blank slide',
        'Advanced 2-Column Text', 'Advanced 3-Column Text', 'Advanced 4-Column Text', 'Advanced Chart Full Width',
        'Advanced Chart 2/3', 'Advanced Chart 2 Column', '1_Standard 1-Column Text',
    ],
    'ICLUSIG': [
        'Title Slide', 'Title_Eyebrow_subhead_Content', '1_Title_Eyebrow_subhead', 'Title',
        'Section Header', 'Title Only', 'Standard 2-Column Text', 'Title Page 1',
        '1_Title Page 2 ', 'Section Title Page', 'Blank slide', 'Advanced 2-Column Text',
        'Advanced 3-Column Text', 'Advanced Chart Full Width', 'Advanced Chart 2/3', 'Advanced Chart 2 Column',
        'Advanced Chart 3 Column', 'Title w/ Subhead & Headline', 'Title and Content', 'Comparison',
        'A - Title and Content', 'Title Slide 1', 'Title Slide 1 ALT', 'Title Slide 2',
        'Section Header 1', 'Section Header 2', 'TOC', 'A - Title Only',
        'B - Title and Two Content', 'B - Title Only', 'C - Title and Two Content', 'C - Title Only',
        'Charts and Tables', 'Thank You 1', 'Thank You 2', '1_Standard 1-Column Text',
    ],
    'ENTYVIO': [
        'Title Page 1', '1_Title Page 2 ', 'Standard 2-Column Text', 'Section Title Page',
        'Blank slide', 'Advanced 3-Column Text', 'Advanced Chart Full Width', 'Advanced Chart 2/3',
        'Advanced Chart 2 Column', 'Advanced Chart 3 Column', 'Title Only', 'Title',
        'Title and Content', 'Title Only', 'Section Divider Dark', 'Section Divider Light',
        '2 Column', '3 Column', 'One Third LHS Dark', 'One Third LHS Light',
        'Half Image Dark', 'Half Image Light', 'Full Image', 'Focus/Callout',
        'Blank', 'End', 'Title Slide', '1_Title & Content Grey',
        '7_Custom Layout', '2_Standard 1-Column Text', 'Section',
    ],
}

# Brand chrome: which elements each brand stamps, and the helper that applies them.
# (Geometry lives in takeda_templates.CHROME_GEOMETRY.)
BRAND_CHROME = {
    'FRUZAQLA': {'banner': 'FRUQ_BANNER (top gradient, full-bleed, thin bottom border line)',
                 'logo':   'FRUQ_LOGO (bottom-right, whitespace-contingent)',
                 'apply':  'apply_fruzaqla_chrome(slide, place_logo="auto")'},
    'ENTYVIO':  {'banner': 'ENTYVIO_BANNER (bottom full-width gradient strip, always applied)',
                 'apply':  'apply_entyvio_chrome(slide)'},
    'ICLUSIG':  {'logo':   'ICLUSIG_LOGO (bottom-right, whitespace-contingent)',
                 'apply':  'apply_iclusig_chrome(slide, place_logo="auto")'},
}

def get_brand(name):
    """Return the brand config dict, or raise with the valid keys."""
    try:
        return BRANDS[name.upper()]
    except KeyError:
        raise KeyError(f"Unknown brand {name!r}. Known: {sorted(BRANDS)}")



# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — EXECUTIVE SUMMARY SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

# --- build_stacked_findings_slide ---
def build_stacked_findings_slide(deck, title, findings, category_label=None, footer_text=None):
    """
    Build one Stacked Findings exec summary slide.
    
    Args:
        deck: Deck instance
        title: Slide title string (e.g. "Executive Summary")
        findings: List of dicts, each with:
            - 'headline': str (bold accent text)
            - 'details': list of str (supporting bullets)
            - 'sub_items': list of str (optional, "Label: detail" format)
        category_label: Optional sidebar label
        footer_text: Optional footer/sources text
    
    Max 4 findings per slide. If more, split across slides.
    """
    slide = deck.add_slide("ONE_COLUMN")
    remove_body_placeholder(slide)
    
    # Title
    deck.set_placeholder_text(slide, 0, title)
    constrain_title_placeholder(slide)
    
    # Background panel (charcoal)
    bg = slide.shapes.add_shape(
        1, Inches(1.174), Inches(1.230), Inches(12.159), Inches(5.612)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(CHARCOAL)
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    
    # Vertical accent line
    line = slide.shapes.add_shape(
        1, Inches(2.498), Inches(1.433), Inches(0.020), Inches(5.200)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
    line.line.fill.background()
    
    # Finding blocks
    current_top = 1.281
    BLOCK_LEFT = 2.765
    BLOCK_WIDTH = 10.233
    GAP = 0.080
    
    for finding in findings[:4]:  # max 4
        headline = finding.get('headline', '')
        details = finding.get('details', [])
        sub_items = finding.get('sub_items', [])
        
        # Estimate height
        line_count = 1 + len(details) + len(sub_items)
        block_height = max(0.70, 0.28 * line_count)
        
        # Create text box (transparent fill)
        txBox = slide.shapes.add_textbox(
            Inches(BLOCK_LEFT), Inches(current_top),
            Inches(BLOCK_WIDTH), Inches(block_height)
        )
        txBox.fill.background()  # transparent
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Headline paragraph
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = headline
        run.font.name = BRAND_FONT
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(ACCENT)
        
        # Detail paragraphs
        for detail in details:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 1
            run = p.add_run()
            run.text = detail
            run.font.name = BRAND_FONT
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.color.rgb = RGBColor.from_string(WHITE)  # white on dark bg
            # Add spacing
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            spcBef = pPr.find(qn('a:spcBef'))
            if spcBef is None:
                spcBef = etree.SubElement(pPr, qn('a:spcBef'))
            spcPts = spcBef.find(qn('a:spcPts'))
            if spcPts is None:
                spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', '200')  # 2pt before
        
        # Sub-item paragraphs
        for sub in sub_items:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            run = p.add_run()
            run.text = sub
            run.font.name = BRAND_FONT
            run.font.size = Pt(12)
            run.font.bold = True  # label portion
            run.font.color.rgb = RGBColor.from_string(WHITE)
        
        current_top += block_height + GAP
    
    # Footer
    if footer_text:
        ft = slide.shapes.add_textbox(
            Inches(-0.067), Inches(6.926), Inches(12.563), Inches(0.343)
        )
        ft.text_frame.word_wrap = True
        p = ft.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = footer_text
        run.font.name = BRAND_FONT
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor.from_string(CHARCOAL)
    
    return slide

# --- build_card_grid_slide ---
def build_card_grid_slide(deck, title, banner_title, cards):
    """
    Build one Card Grid exec summary slide.
    
    Args:
        deck: Deck instance
        title: Slide title (goes in placeholder idx=0)
        banner_title: Centered banner text (e.g. "Key Recommendations")
        cards: List of 3-4 dicts, each with:
            - 'label': str (category label, 1-2 words)
            - 'description': str or list of str (bullet points)
    """
    slide = deck.add_slide("ONE_COLUMN")
    remove_body_placeholder(slide)
    
    # Title
    deck.set_placeholder_text(slide, 0, title)
    constrain_title_placeholder(slide)
    
    n = len(cards)
    assert 2 <= n <= 4, "Card grid supports 2-4 cards"
    
    # Column geometry
    if n == 4:
        col_lefts   = [0.668, 3.789, 6.910, 9.900]
        circ_lefts  = [1.311, 4.431, 7.552, 10.669]
        desc_lefts  = [0.551, 3.716, 6.805, 9.978]
        col_widths  = [2.602, 2.602, 2.602, 2.862]
        desc_widths = [2.700, 2.752, 2.820, 2.862]
    elif n == 3:
        col_lefts   = [1.100, 4.600, 8.100]
        circ_lefts  = [1.842, 5.342, 8.842]
        desc_lefts  = [1.000, 4.500, 8.000]
        col_widths  = [2.800, 2.800, 2.800]
        desc_widths = [2.900, 2.900, 2.900]
    else:  # n == 2
        col_lefts   = [2.000, 7.200]
        circ_lefts  = [2.842, 8.042]
        desc_lefts  = [1.900, 7.100]
        col_widths  = [3.200, 3.200]
        desc_widths = [3.300, 3.300]
    
    # Banner bar
    if banner_title:
        banner_w = min(6.789, n * 2.5)
        banner_l = (13.333 - banner_w) / 2  # centered
        banner = slide.shapes.add_shape(
            1, Inches(banner_l), Inches(1.151), Inches(banner_w), Inches(0.303)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor.from_string(DARK_ACC)
        banner.line.fill.background()
        banner.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = banner.text_frame.paragraphs[0].add_run()
        run.text = banner_title
        run.font.name = BRAND_FONT
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(WHITE)
    
    for i, card in enumerate(cards):
        # Icon circle
        from pptx.enum.shapes import MSO_SHAPE
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(circ_lefts[i]), Inches(2.398),
            Inches(1.316), Inches(1.316)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor.from_string(ACCENT.replace('B1254A', 'E1242A'))
        circle.line.fill.background()
        
        # Category label bar
        label_text = card.get('label', '')
        label_lines = 1 if len(label_text) <= 25 else 2
        label_h = 0.303 if label_lines == 1 else 0.606
        label = slide.shapes.add_shape(
            1,
            Inches(col_lefts[i]), Inches(3.937),
            Inches(col_widths[i]), Inches(label_h)
        )
        label.fill.solid()
        label.fill.fore_color.rgb = RGBColor.from_string(DARK_ACC)
        label.line.fill.background()
        label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        label.text_frame.word_wrap = True
        run = label.text_frame.paragraphs[0].add_run()
        run.text = label_text
        run.font.name = BRAND_FONT
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(WHITE)
        
        # Description box
        desc_top = 3.937 + label_h + 0.10
        desc_items = card.get('description', '')
        if isinstance(desc_items, str):
            desc_items = [desc_items]
        
        desc_h = max(1.35, 0.45 * len(desc_items))
        desc = slide.shapes.add_textbox(
            Inches(desc_lefts[i]), Inches(desc_top),
            Inches(desc_widths[i]), Inches(desc_h)
        )
        desc.text_frame.word_wrap = True
        for j, item in enumerate(desc_items):
            if j == 0:
                p = desc.text_frame.paragraphs[0]
            else:
                p = desc.text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = item
            run.font.name = BRAND_FONT
            run.font.size = Pt(16)
            run.font.bold = False
            run.font.color.rgb = RGBColor.from_string(CHARCOAL)
    
    return slide

# --- chunk_findings (from procedure block) ---
def chunk_findings(findings, max_per_slide=4):
    """Split a flat list of findings into slide-sized chunks."""
    return [findings[i:i + max_per_slide]
            for i in range(0, len(findings), max_per_slide)]


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — QA GATE  (also importable: from takeda import validate)
# Run: python takeda.py output.pptx
# ═══════════════════════════════════════════════════════════════════════════════
import sys, re
from pptx import Presentation
from pptx.util import Emu
try:
    from takeda_deck import TAKEDA_PALETTE
except Exception:
    # fallback mirrors takeda_deck; keep in sync via that module, not here
    TAKEDA_PALETTE = {'E1242A','891515','34373F','A1B1C3','D1D8E0','EDF2F3',
                      'A1A4AC','EDF2F4','FFFFFF','000000'}

EMU = 914400
PROMPT_RE = re.compile(r'click to add|insert photo|lorem|todo|page\s+\d+\s+of\s+\d+', re.I)

def _iter_runs(prs):
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        yield si, sh, r

def validate(path, brand=None, mode="A"):
    """QA gate. With brand=None, validates against the corporate TAKEDA_PALETTE.
    With brand in {FRUZAQLA,ICLUSIG,ENTYVIO}, validates colors against that brand's
    full harvested palette (BRAND_FULL_PALETTE) — corporate neutrals are always allowed.

    mode="A" (default — authoring a NEW deck) enforces the 12pt-minimum and
    left-accent-zone rules as HARD failures.
    mode="B" (reformatting an EXISTING deck) preserves the source's original font
    sizes and geometry, so those two checks become non-blocking WARNINGS: a source
    legitimately carries kicker tags / footnotes / dense table cells under 12pt and
    full-width chrome / off-canvas bleeds in the left zone, and bumping them would
    distort the source. Off-palette colours and leftover prompt text stay HARD
    failures in BOTH modes (they are real defects regardless of mode).

    Returns (fails, warnings) — fails block the deck (exit 1); warnings are printed
    for awareness but do not block.
    """
    prs = Presentation(path)
    reformat = str(mode).upper() == "B"
    fails = []
    warnings = []

    # Resolve the acceptable palette for this deck.
    if brand and brand.upper() in BRAND_FULL_PALETTE:
        allowed = set(BRAND_FULL_PALETTE[brand.upper()]) | TAKEDA_PALETTE
    else:
        allowed = set(TAKEDA_PALETTE)

    # 1. Off-palette colors (explicit srgbClr in slide XML)
    def _sat(hx):
        try: r, g, b = (int(hx[i:i+2], 16) for i in (0, 2, 4))
        except Exception: return 1.0
        mx, mn = max(r, g, b), min(r, g, b)
        return 0.0 if mx == 0 else (mx - mn) / mx
    off = set()
    off_neutral = set()
    for si, s in enumerate(prs.slides):
        for el in s.shapes._spTree.iter():
            val = el.get('val')
            if val and el.tag.endswith('}srgbClr'):
                hx = val.upper()
                if hx not in allowed:
                    # Mode B tolerates PRESERVED SOURCE NEUTRALS (near-gray, low saturation)
                    # — e.g. a data-viz gray ramp kept to match the source — as warnings.
                    # Chromatic off-brand colours (e.g. leftover orange) still HARD-fail.
                    if reformat and _sat(hx) < 0.12:
                        off_neutral.add((si+1, hx))
                    else:
                        off.add((si+1, hx))
    if off:
        fails.append(f"Off-palette colors{f' (brand={brand})' if brand else ''}: {sorted(off)[:12]}")
    if off_neutral:
        warnings.append(f"Off-palette NEUTRALs (preserved source grays, mode B): {sorted(off_neutral)[:12]}")

    # 2. Body text below 12pt (excluding footnote zone, bottom 0.4")
    small = []
    for si, sh, r in _iter_runs(prs):
        if r.font.size and r.font.size.pt < 12:
            top_in = (sh.top or 0) / EMU
            if top_in < 7.0:  # not a footnote
                small.append((si+1, round(r.font.size.pt,1), (r.text or '')[:20]))
    if small:
        # Mode A: hard fail (authoring guardrail). Mode B: warn only — the source's
        # original sizes are preserved on purpose (kickers/footnotes/table cells).
        (warnings if reformat else fails).append(f"Sub-12pt text: {small[:8]}")

    # 3. Placeholder prompt / leftover text
    leftovers = []
    for si, sh, r in _iter_runs(prs):
        if PROMPT_RE.search(r.text or ''):
            leftovers.append((si+1, r.text[:30]))
    if leftovers:
        fails.append(f"Leftover/prompt text: {leftovers[:8]}")

    # 4. Content intruding into left-accent exclusion zone (L < 0.6", below title)
    #    Skip layout-inherited placeholders — they ARE the chrome and legitimately sit there.
    intrude = []
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.left is None or sh.top is None:
                continue
            if getattr(sh, "is_placeholder", False):
                continue
            if sh.top/EMU >= 1.10 and sh.left/EMU < 0.60:
                intrude.append((si+1, round(sh.left/EMU,2)))
    if intrude:
        # Mode A: hard fail. Mode B: warn only — full-width footer/rules/charts and
        # off-canvas design bleeds start at the margin and trip this by geometry, not error.
        (warnings if reformat else fails).append(f"Shapes in left-accent zone (L<0.6\"): {intrude[:8]}")

    return fails, warnings

if __name__ == "__main__":
    import sys
    args = sys.argv[1:]

    if args and args[0] == "--regen-layouts":
        if len(args) < 2:
            print("usage: python takeda.py --regen-layouts <template.potx>"); sys.exit(2)
        import json as _json
        d = Deck(args[1])
        reg = {}
        names = [L.name for L in d.prs.slide_masters[0].slide_layouts]
        for key,(idx,name) in LAYOUTS.items():
            if idx < len(names) and names[idx]==name:
                reg[key] = {"name":name,"master":0,"layout":idx}
        for key,(mi,li) in PATTERN_LAYOUTS.items():
            reg[key] = {"name":"pattern","master":mi,"layout":li}
        print(_json.dumps(reg, indent=2))
        sys.exit(0)

    if not args:
        print("usage: python takeda.py <deck.pptx> [--brand FRUZAQLA|ICLUSIG|ENTYVIO] [--mode A|B]\n"
              "       --mode A (default) = authoring a new deck (strict).\n"
              "       --mode B / --reformat = reformatting an existing deck (12pt-min and\n"
              "         left-accent-zone become non-blocking warnings; source sizes preserved).")
        sys.exit(2)
    brand = None
    if "--brand" in args:
        i = args.index("--brand")
        brand = args[i+1] if i+1 < len(args) else None
        args = args[:i] + args[i+2:]
    mode = "A"
    if "--reformat" in args:
        mode = "B"; args = [a for a in args if a != "--reformat"]
    if "--mode" in args:
        i = args.index("--mode")
        mode = (args[i+1] if i+1 < len(args) else "A").upper()
        args = args[:i] + args[i+2:]

    fails, warnings = validate(args[0], brand=brand, mode=mode)
    if warnings:
        print(f"\u26a0\ufe0f  QA warnings (mode B — reformatting; non-blocking):")
        for w in warnings: print("  -", w)
    if fails:
        print("\u274c QA FAILED — do not share:")
        for f in fails: print("  -", f)
        sys.exit(1)
    tag = " (warnings only)" if warnings else ""
    print(f"\u2705 QA passed{tag} — safe to present_files."); sys.exit(0)
